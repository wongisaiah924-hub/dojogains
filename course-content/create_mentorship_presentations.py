#!/usr/bin/env python3
"""
Dojo Gains - Mentorship Presentations
Creates all 9 mentorship module presentations
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
DARK_BG = RGBColor(15, 15, 26)
RED_ACCENT = RGBColor(230, 57, 70)
GREEN = RGBColor(34, 197, 94)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(156, 163, 175)
BLUE = RGBColor(59, 130, 246)
YELLOW = RGBColor(234, 179, 8)

SLIDE_W = 10
SLIDE_H = 7.5

def set_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK_BG

def text(slide, txt, top, left, width, size=18, color=WHITE, bold=False, align="center"):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = PP_ALIGN.CENTER if align == "center" else PP_ALIGN.LEFT
    return box

def title(slide, txt, top=0.4):
    text(slide, txt, top, 0.5, 9, size=40, color=WHITE, bold=True)

def bullets(slide, points, top, size=20):
    box = slide.shapes.add_textbox(Inches(0.75), Inches(top), Inches(8.5), Inches(5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, pt in enumerate(points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•   {pt}"
        p.font.size = Pt(size)
        p.font.color.rgb = WHITE
        p.space_before = Pt(6)
        p.space_after = Pt(10)

def footer(slide, txt):
    text(slide, txt, 6.9, 0.5, 9, size=14, color=GRAY)

def hline(slide, top, left, width, color=RED_ACCENT):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.02))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()

# ============================================
# 1. BUILD A 7 FIGURE MINDSET
# ============================================
def create_mindset_ppt():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    text(slide, "BUILD A", 1.8, 0.5, 9, size=28, color=GRAY)
    text(slide, "7 FIGURE MINDSET", 2.5, 0.5, 9, size=48, color=WHITE, bold=True)
    hline(slide, 3.4, 3.5, 3)
    text(slide, "DOJO GAINS - Mentorship", 3.9, 0.5, 9, size=20, color=RED_ACCENT, bold=True)
    text(slide, "The Mental Game of Trading", 4.4, 0.5, 9, size=18, color=GRAY)
    
    # Slide 2: Why Mindset Matters
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Why Mindset Matters")
    bullets(slide, [
        "Trading is 80% psychology, 20% strategy",
        "Your emotions are your biggest enemy",
        "Most traders fail not from bad strategies, but bad discipline",
        "The market is designed to exploit human psychology",
        "7-figure traders think differently than retail traders"
    ], 1.2, size=22)
    footer(slide, "Master your mind, master the markets")
    
    # Slide 3: The Losing Trader Mindset
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "The Losing Trader Mindset")
    bullets(slide, [
        "Trades emotionally - revenge trades after losses",
        "Thinks every trade must be a winner",
        "Moves stop losses to avoid taking losses",
        "Over-leverages to 'make back' losses",
        "Blames the market, brokers, or others",
        "Chases trades and enters without a plan",
        "Can't walk away from the screen"
    ], 1.2, size=20)
    footer(slide, "Recognize these patterns? Time to change.")
    
    # Slide 4: The Winning Trader Mindset
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "The Winning Trader Mindset")
    bullets(slide, [
        "Accepts losses as part of the business",
        "Focuses on process, not individual outcomes",
        "Follows the plan no matter what",
        "Takes responsibility for every trade",
        "Stays patient - waits for A+ setups",
        "Protects capital above all else",
        "Knows when to step away"
    ], 1.2, size=20)
    footer(slide, "This is the mindset that builds wealth")
    
    # Slide 5: Think in Probabilities
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Think in Probabilities")
    bullets(slide, [
        "No single trade matters - it's the SERIES that counts",
        "A 60% win rate means 40 losses per 100 trades",
        "You can lose 5 trades in a row and still be profitable",
        "Each trade is just one sample in a large dataset",
        "Remove emotion by thinking statistically",
        "Focus on expected value, not individual results"
    ], 1.2, size=22)
    footer(slide, "Think like a casino, not a gambler")
    
    # Slide 6: The Process Over Outcome
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Process Over Outcome")
    bullets(slide, [
        "A good trade can lose money",
        "A bad trade can make money",
        "Judge yourself on EXECUTION, not P&L",
        "Did you follow your rules? That's a good trade.",
        "Did you break your rules? That's a bad trade - even if it won.",
        "Consistent process = consistent results over time"
    ], 1.2, size=22)
    footer(slide, "The process is the prize")
    
    # Slide 7: Handling Losses
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Handling Losses Like a Pro")
    bullets(slide, [
        "Losses are tuition - you're paying to learn",
        "Every loss teaches you something",
        "Take the loss, journal it, move on",
        "NEVER revenge trade - walk away if emotional",
        "Set a daily loss limit (e.g., 2-3% max)",
        "Hit the limit? You're done for the day. No exceptions."
    ], 1.2, size=22)
    footer(slide, "How you handle losses defines your career")
    
    # Slide 8: Building Discipline
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Building Unshakeable Discipline")
    bullets(slide, [
        "Create a trading plan and follow it religiously",
        "Use a pre-trade checklist before every entry",
        "Set rules for max trades per day",
        "Take breaks - trading is mentally exhausting",
        "Journal every trade - winners AND losers",
        "Review weekly - what did you do well? What to improve?"
    ], 1.2, size=22)
    footer(slide, "Discipline is freedom")
    
    # Slide 9: The Long Game
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Playing the Long Game")
    bullets(slide, [
        "7 figures isn't built in a day, month, or even year",
        "Focus on consistent small gains, not home runs",
        "Compound your account - 1% per day = 1,200%+ per year",
        "Protect your capital during learning phase",
        "The goal is to STAY in the game long enough to master it",
        "Most quit right before they would have succeeded"
    ], 1.2, size=22)
    footer(slide, "The tortoise beats the hare")
    
    # Slide 10: Daily Habits of 7-Figure Traders
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Daily Habits of 7-Figure Traders")
    bullets(slide, [
        "Morning routine: Review markets, mark levels, set plan",
        "Only trade during their best hours (Killzones)",
        "Take breaks every 90 minutes",
        "Exercise and stay physically healthy",
        "Journal every trade same day",
        "End of day review: What worked? What didn't?",
        "Continuous education - never stop learning"
    ], 1.2, size=20)
    
    # Slide 11: Affirmations
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Trader Affirmations")
    text(slide, "I am a disciplined trader.", 1.4, 0.5, 9, size=24, color=GREEN)
    text(slide, "I follow my plan no matter what.", 2.0, 0.5, 9, size=24, color=GREEN)
    text(slide, "Losses are part of the game.", 2.6, 0.5, 9, size=24, color=GREEN)
    text(slide, "I think in probabilities, not certainties.", 3.2, 0.5, 9, size=24, color=GREEN)
    text(slide, "I protect my capital above all.", 3.8, 0.5, 9, size=24, color=GREEN)
    text(slide, "I am patient. The setup will come.", 4.4, 0.5, 9, size=24, color=GREEN)
    text(slide, "I am building generational wealth.", 5.0, 0.5, 9, size=24, color=GREEN)
    footer(slide, "Read these every morning before trading")
    
    # Slide 12: Key Takeaways
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Key Takeaways")
    bullets(slide, [
        "Mindset separates winners from losers",
        "Think in probabilities, not certainties",
        "Process over outcome - always",
        "Losses are tuition, not failure",
        "Discipline is your superpower",
        "Play the long game - compound wins"
    ], 1.2, size=22)
    text(slide, "Your mind is your edge", 5.5, 0.5, 9, size=22, color=RED_ACCENT, bold=True)
    
    prs.save("/root/clawd/dojo-gains/course-content/M01_7_Figure_Mindset.pptx")
    print("Saved: M01_7_Figure_Mindset.pptx")

# ============================================
# 2. FAIR VALUE GAPS AND IFVGs
# ============================================
def create_fvg_deep_dive_ppt():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    text(slide, "FAIR VALUE GAPS", 2.0, 0.5, 9, size=48, color=WHITE, bold=True)
    text(slide, "& Inverse FVGs (IFVGs)", 2.9, 0.5, 9, size=28, color=GRAY)
    hline(slide, 3.7, 3.5, 3)
    text(slide, "DOJO GAINS - Mentorship", 4.2, 0.5, 9, size=20, color=RED_ACCENT, bold=True)
    
    # Slide 2: FVG Recap
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "FVG Quick Recap")
    bullets(slide, [
        "Fair Value Gap = Price imbalance on the chart",
        "Created when price moves aggressively",
        "3-candle formation with a gap between 1st and 3rd candle",
        "Bullish FVG: Candle 3 LOW > Candle 1 HIGH",
        "Bearish FVG: Candle 3 HIGH < Candle 1 LOW",
        "Price tends to return and 'fill' these gaps"
    ], 1.2, size=22)
    footer(slide, "Imbalances seek balance")
    
    # Slide 3: FVG Types
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Types of Fair Value Gaps")
    text(slide, "STANDARD FVG", 1.3, 0.5, 4.5, size=22, color=GREEN, bold=True)
    text(slide, "Gap between wicks of candle 1 & 3", 1.7, 0.5, 4.5, size=16, color=GRAY)
    text(slide, "Most common type", 2.0, 0.5, 4.5, size=16, color=GRAY)
    
    text(slide, "CONSEQUENT ENCROACHMENT (CE)", 2.7, 0.5, 9, size=22, color=BLUE, bold=True)
    text(slide, "The 50% midpoint of the FVG - key entry level", 3.1, 0.5, 9, size=16, color=GRAY)
    
    text(slide, "INVERSE FVG (IFVG)", 3.8, 0.5, 4.5, size=22, color=RED_ACCENT, bold=True)
    text(slide, "When an FVG gets violated", 4.2, 0.5, 4.5, size=16, color=GRAY)
    text(slide, "and becomes the opposite zone", 4.5, 0.5, 4.5, size=16, color=GRAY)
    
    text(slide, "BREAKER + FVG", 3.8, 5, 4.5, size=22, color=YELLOW, bold=True)
    text(slide, "FVG that forms after BOS", 4.2, 5, 4.5, size=16, color=GRAY)
    text(slide, "Higher probability zone", 4.5, 5, 4.5, size=16, color=GRAY)
    
    # Slide 4: What is an IFVG?
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "What is an Inverse FVG (IFVG)?")
    bullets(slide, [
        "An IFVG forms when a regular FVG gets VIOLATED",
        "Price trades through the FVG instead of respecting it",
        "The violated FVG flips to become support/resistance",
        "Bullish FVG gets violated = Now acts as RESISTANCE",
        "Bearish FVG gets violated = Now acts as SUPPORT",
        "Shows a shift in market sentiment"
    ], 1.2, size=22)
    footer(slide, "When balance becomes imbalance again")
    
    # Slide 5: IFVG Formation
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "How IFVG Forms")
    bullets(slide, [
        "1. A standard FVG forms (bullish or bearish)",
        "2. Price returns to the FVG as expected",
        "3. BUT instead of bouncing, price BREAKS through",
        "4. The FVG is now 'inversed' - flipped",
        "5. Old support becomes resistance (and vice versa)",
        "6. Look for entries at the IFVG on retest"
    ], 1.2, size=22)
    footer(slide, "Failed FVG = IFVG opportunity")
    
    # Slide 6: Trading IFVGs
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "How to Trade IFVGs")
    bullets(slide, [
        "Identify a violated FVG (price broke through)",
        "Wait for price to retrace back to the IFVG",
        "Enter in the NEW direction (opposite of original FVG)",
        "Bullish FVG violated = SHORT at IFVG",
        "Bearish FVG violated = LONG at IFVG",
        "Stop beyond IFVG, target next liquidity"
    ], 1.2, size=22)
    footer(slide, "The flip is the setup")
    
    # Slide 7: Consequent Encroachment
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Consequent Encroachment (CE)")
    bullets(slide, [
        "CE = The 50% midpoint of any FVG",
        "This is the EQUILIBRIUM of the imbalance",
        "Price often reacts precisely at the CE level",
        "Use CE as your primary entry point",
        "More precise than trading the entire FVG zone",
        "Mark the CE on every FVG you identify"
    ], 1.2, size=22)
    footer(slide, "50% is the sweet spot")
    
    # Slide 8: FVG Quality Filters
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "FVG Quality Filters")
    bullets(slide, [
        "NOT all FVGs are created equal. Filter for:",
        "1. FVG in line with higher timeframe trend",
        "2. FVG formed during Killzone hours",
        "3. FVG near a key level (support/resistance)",
        "4. FVG after a liquidity sweep",
        "5. FVG with strong displacement candle",
        "6. First FVG in a new move (not 3rd or 4th)"
    ], 1.2, size=20)
    footer(slide, "Quality over quantity")
    
    # Slide 9: FVG + Order Block
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "FVG + Order Block = A+ Setup")
    bullets(slide, [
        "Order Block = Last candle before displacement",
        "FVG = The gap created by displacement",
        "When they overlap = Highest probability zone",
        "This is where institutions placed their orders",
        "Look for OB + FVG confluence for entries",
        "Even better: OB + FVG + Key Level"
    ], 1.2, size=22)
    footer(slide, "Confluence is king")
    
    # Slide 10: Common FVG Mistakes
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Common FVG Mistakes")
    bullets(slide, [
        "Trading every FVG regardless of context",
        "Not waiting for price to reach the FVG",
        "Entering before confirmation at the FVG",
        "Trading FVGs against the trend",
        "Ignoring violated FVGs (missing IFVG trades)",
        "Using FVGs on very low timeframes (noise)"
    ], 1.2, size=22)
    footer(slide, "Avoid these and level up")
    
    # Slide 11: Key Takeaways
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Key Takeaways")
    bullets(slide, [
        "FVGs show where price moved too fast",
        "CE (50%) is your primary entry level",
        "IFVG = Violated FVG that flipped direction",
        "FVG + Order Block = High probability setup",
        "Filter FVGs by trend, time, and context",
        "Quality FVGs with confluence = Best trades"
    ], 1.2, size=22)
    text(slide, "Master FVGs, master entries", 5.5, 0.5, 9, size=22, color=RED_ACCENT, bold=True)
    
    prs.save("/root/clawd/dojo-gains/course-content/M02_FVG_and_IFVG.pptx")
    print("Saved: M02_FVG_and_IFVG.pptx")

# ============================================
# 3. DIVERGENCES
# ============================================
def create_divergences_ppt():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    text(slide, "DIVERGENCES", 2.2, 0.5, 9, size=52, color=WHITE, bold=True)
    text(slide, "When Price and Momentum Disagree", 3.1, 0.5, 9, size=26, color=GRAY)
    hline(slide, 3.9, 3.5, 3)
    text(slide, "DOJO GAINS - Mentorship", 4.4, 0.5, 9, size=20, color=RED_ACCENT, bold=True)
    
    # Slide 2: What is Divergence?
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "What is Divergence?")
    bullets(slide, [
        "Divergence = Price and indicator moving in opposite directions",
        "Shows weakening momentum behind a move",
        "Early warning sign of potential reversal",
        "Common indicators: RSI, MACD, Stochastic",
        "One of the most powerful reversal signals",
        "Works on all timeframes and markets"
    ], 1.2, size=22)
    footer(slide, "Divergence reveals what price hides")
    
    # Slide 3: Types of Divergence
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Types of Divergence")
    text(slide, "REGULAR DIVERGENCE", 1.3, 0.5, 4.5, size=22, color=GREEN, bold=True)
    text(slide, "Signals potential REVERSAL", 1.7, 0.5, 4.5, size=16, color=GRAY)
    text(slide, "Trend is weakening", 2.0, 0.5, 4.5, size=16, color=GRAY)
    
    text(slide, "HIDDEN DIVERGENCE", 1.3, 5, 4.5, size=22, color=BLUE, bold=True)
    text(slide, "Signals potential CONTINUATION", 1.7, 5, 4.5, size=16, color=GRAY)
    text(slide, "Trend is strong, pullback ending", 2.0, 5, 4.5, size=16, color=GRAY)
    
    bullets(slide, [
        "Regular = Counter-trend signal (reversal)",
        "Hidden = With-trend signal (continuation)",
        "Both are powerful when used correctly"
    ], 3.0, size=20)
    
    # Slide 4: Regular Bullish Divergence
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Regular Bullish Divergence")
    bullets(slide, [
        "Price makes LOWER LOW",
        "Indicator makes HIGHER LOW",
        "Meaning: Selling pressure is weakening",
        "Buyers are stepping in despite lower prices",
        "Signal: Potential bullish reversal coming",
        "Best at: Key support levels, after extended downtrend"
    ], 1.2, size=22)
    text(slide, "Price: Lower Low | Indicator: Higher Low = BUY SIGNAL", 5.5, 0.5, 9, size=18, color=GREEN, bold=True)
    
    # Slide 5: Regular Bearish Divergence
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Regular Bearish Divergence")
    bullets(slide, [
        "Price makes HIGHER HIGH",
        "Indicator makes LOWER HIGH",
        "Meaning: Buying pressure is weakening",
        "Sellers are stepping in despite higher prices",
        "Signal: Potential bearish reversal coming",
        "Best at: Key resistance levels, after extended uptrend"
    ], 1.2, size=22)
    text(slide, "Price: Higher High | Indicator: Lower High = SELL SIGNAL", 5.5, 0.5, 9, size=18, color=RED_ACCENT, bold=True)
    
    # Slide 6: Hidden Bullish Divergence
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Hidden Bullish Divergence")
    bullets(slide, [
        "Price makes HIGHER LOW (uptrend pullback)",
        "Indicator makes LOWER LOW",
        "Meaning: Pullback is losing momentum",
        "Trend is still strong, continuation likely",
        "Signal: Uptrend will resume",
        "Best for: Trend continuation entries"
    ], 1.2, size=22)
    text(slide, "Uptrend pullback ending = BUY the dip", 5.5, 0.5, 9, size=18, color=GREEN, bold=True)
    
    # Slide 7: Hidden Bearish Divergence
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Hidden Bearish Divergence")
    bullets(slide, [
        "Price makes LOWER HIGH (downtrend rally)",
        "Indicator makes HIGHER HIGH",
        "Meaning: Rally is losing momentum",
        "Trend is still strong, continuation likely",
        "Signal: Downtrend will resume",
        "Best for: Trend continuation entries"
    ], 1.2, size=22)
    text(slide, "Downtrend rally ending = SELL the rip", 5.5, 0.5, 9, size=18, color=RED_ACCENT, bold=True)
    
    # Slide 8: Best Indicators for Divergence
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Best Indicators for Divergence")
    text(slide, "RSI (Relative Strength Index)", 1.3, 0.5, 9, size=22, color=GREEN, bold=True)
    text(slide, "Most popular, works on all timeframes, clear signals", 1.7, 0.5, 9, size=16, color=GRAY)
    
    text(slide, "MACD", 2.4, 0.5, 9, size=22, color=BLUE, bold=True)
    text(slide, "Good for trend strength, histogram divergence is powerful", 2.8, 0.5, 9, size=16, color=GRAY)
    
    text(slide, "Stochastic", 3.5, 0.5, 9, size=22, color=YELLOW, bold=True)
    text(slide, "Fast signals, best for ranging markets", 3.9, 0.5, 9, size=16, color=GRAY)
    
    text(slide, "OBV (On-Balance Volume)", 4.6, 0.5, 9, size=22, color=WHITE, bold=True)
    text(slide, "Volume-based divergence, shows institutional activity", 5.0, 0.5, 9, size=16, color=GRAY)
    
    # Slide 9: Trading Divergence
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "How to Trade Divergence")
    bullets(slide, [
        "1. Identify the divergence (price vs indicator)",
        "2. Confirm with key level (S/R, FVG, order block)",
        "3. Wait for candlestick confirmation",
        "4. Enter after confirmation candle closes",
        "5. Stop beyond the swing high/low",
        "6. Target next key level or liquidity"
    ], 1.2, size=22)
    footer(slide, "Divergence + Confluence = High probability")
    
    # Slide 10: Divergence Mistakes
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Common Divergence Mistakes")
    bullets(slide, [
        "Trading divergence alone without confluence",
        "Entering before confirmation",
        "Fighting strong trends with divergence",
        "Using divergence on very low timeframes",
        "Ignoring the bigger picture (HTF trend)",
        "Not waiting for the second swing to complete"
    ], 1.2, size=22)
    footer(slide, "Divergence is a warning, not a trigger")
    
    # Slide 11: Key Takeaways
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Key Takeaways")
    bullets(slide, [
        "Regular Divergence = Reversal signal",
        "Hidden Divergence = Continuation signal",
        "Bullish: Price LL, Indicator HL (or vice versa)",
        "Bearish: Price HH, Indicator LH (or vice versa)",
        "RSI is the most reliable for divergence",
        "Always combine with key levels for confirmation"
    ], 1.2, size=22)
    text(slide, "Divergence + Level + Confirmation = Trade", 5.5, 0.5, 9, size=20, color=RED_ACCENT, bold=True)
    
    prs.save("/root/clawd/dojo-gains/course-content/M03_Divergences.pptx")
    print("Saved: M03_Divergences.pptx")

# ============================================
# 4. ANATOMY OF A CANDLE (Deep Dive)
# ============================================
def create_candle_deep_dive_ppt():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    text(slide, "ANATOMY OF A CANDLE", 2.2, 0.5, 9, size=44, color=WHITE, bold=True)
    text(slide, "Deep Dive into Price Action", 3.1, 0.5, 9, size=26, color=GRAY)
    hline(slide, 3.9, 3.5, 3)
    text(slide, "DOJO GAINS - Mentorship", 4.4, 0.5, 9, size=20, color=RED_ACCENT, bold=True)
    
    # Slide 2: Beyond the Basics
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Beyond the Basics")
    bullets(slide, [
        "You know OHLC - but what does it MEAN?",
        "Every candle tells a story of the battle",
        "Understanding candles = reading the market's mind",
        "This session: How to decode institutional activity",
        "Learn to see what retail traders miss"
    ], 1.2, size=22)
    footer(slide, "The candle speaks - learn to listen")
    
    # Slide 3: The Four Elements
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "The Four Elements")
    text(slide, "OPEN", 1.3, 0.5, 4.5, size=24, color=BLUE, bold=True)
    text(slide, "Where the battle began", 1.7, 0.5, 4.5, size=16, color=GRAY)
    
    text(slide, "HIGH", 2.4, 0.5, 4.5, size=24, color=GREEN, bold=True)
    text(slide, "Maximum buyer reach", 2.8, 0.5, 4.5, size=16, color=GRAY)
    
    text(slide, "LOW", 1.3, 5, 4.5, size=24, color=RED_ACCENT, bold=True)
    text(slide, "Maximum seller reach", 1.7, 5, 4.5, size=16, color=GRAY)
    
    text(slide, "CLOSE", 2.4, 5, 4.5, size=24, color=YELLOW, bold=True)
    text(slide, "Who won the battle", 2.8, 5, 4.5, size=16, color=GRAY)
    
    text(slide, "Close relative to Open = Battle outcome", 3.8, 0.5, 9, size=20, color=WHITE, bold=True)
    text(slide, "Wicks = Rejected price levels", 4.3, 0.5, 9, size=20, color=WHITE, bold=True)
    
    # Slide 4: Body Analysis
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Reading the Body")
    bullets(slide, [
        "LARGE BODY = Strong conviction, one side dominated",
        "SMALL BODY = Indecision, balanced battle",
        "Close near HIGH = Buyers in control at close",
        "Close near LOW = Sellers in control at close",
        "Body as % of range = Conviction meter"
    ], 1.2, size=22)
    text(slide, "Body size = Commitment level", 5.0, 0.5, 9, size=20, color=GREEN, bold=True)
    
    # Slide 5: Wick Analysis
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Reading the Wicks")
    bullets(slide, [
        "UPPER WICK = Buyers pushed up but got rejected",
        "LOWER WICK = Sellers pushed down but got rejected",
        "Long wick = Strong rejection at that level",
        "No wick = Full conviction, no rejection",
        "Wicks show where stops were hunted"
    ], 1.2, size=22)
    text(slide, "Wicks = Rejection + Stop hunts", 5.0, 0.5, 9, size=20, color=RED_ACCENT, bold=True)
    
    # Slide 6: Candle Context
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Context is Everything")
    bullets(slide, [
        "Same candle pattern = Different meaning in different contexts",
        "WHERE the candle forms matters more than WHAT it looks like",
        "Hammer at support = Strong buy signal",
        "Hammer in middle of range = Weak signal",
        "Always ask: Where is this candle forming?",
        "Key levels + Candle pattern = High probability"
    ], 1.2, size=22)
    footer(slide, "Location, location, location")
    
    # Slide 7: Institutional Candles
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Spotting Institutional Candles")
    bullets(slide, [
        "DISPLACEMENT CANDLE: Large body, small wicks",
        "Shows aggressive institutional buying/selling",
        "Creates FVGs and breaks structure",
        "Usually followed by retracement (FVG fill)",
        "These are the candles that matter most",
        "Enter on retracement to displacement zone"
    ], 1.2, size=22)
    footer(slide, "Big candle = Big money moved")
    
    # Slide 8: Rejection Candles
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Rejection Candles")
    bullets(slide, [
        "PIN BAR: Long wick, small body at one end",
        "Shows strong rejection of a price level",
        "The wick is the 'lie' - price couldn't hold there",
        "The body is the 'truth' - where price actually closed",
        "Trade in direction of the body, not the wick",
        "Best when formed at key S/R levels"
    ], 1.2, size=22)
    footer(slide, "Rejection = Opportunity")
    
    # Slide 9: Engulfing Analysis
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Engulfing Candle Deep Dive")
    bullets(slide, [
        "True engulfing: Body engulfs previous body",
        "Even stronger: Body engulfs entire previous candle",
        "Shows complete shift in control",
        "First candle = Old sentiment",
        "Second candle = New sentiment taking over",
        "Volume should increase on engulfing candle"
    ], 1.2, size=22)
    footer(slide, "Engulfing = Power shift")
    
    # Slide 10: Reading Candle Sequences
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Reading Candle Sequences")
    bullets(slide, [
        "Single candles are clues, sequences are stories",
        "3 consecutive large candles = Strong trend",
        "Shrinking candles = Momentum fading",
        "Doji after trend = Indecision, possible reversal",
        "Look for the pattern in the sequence",
        "Each candle adds to the narrative"
    ], 1.2, size=22)
    footer(slide, "Read the story, not just the words")
    
    # Slide 11: Key Takeaways
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Key Takeaways")
    bullets(slide, [
        "Body = Who won | Wicks = Who got rejected",
        "Large body = Conviction | Small body = Indecision",
        "Context matters more than pattern",
        "Displacement candles show institutional activity",
        "Rejection wicks at key levels = Opportunity",
        "Read sequences, not just individual candles"
    ], 1.2, size=22)
    text(slide, "The candle tells the truth", 5.5, 0.5, 9, size=22, color=RED_ACCENT, bold=True)
    
    prs.save("/root/clawd/dojo-gains/course-content/M04_Candle_Anatomy_DeepDive.pptx")
    print("Saved: M04_Candle_Anatomy_DeepDive.pptx")

# ============================================
# 5. LIQUIDITY (Deep Dive)
# ============================================
def create_liquidity_deep_dive_ppt():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    text(slide, "LIQUIDITY", 2.0, 0.5, 9, size=52, color=WHITE, bold=True)
    text(slide, "Advanced Concepts & Trading", 2.9, 0.5, 9, size=26, color=GRAY)
    hline(slide, 3.7, 3.5, 3)
    text(slide, "DOJO GAINS - Mentorship", 4.2, 0.5, 9, size=20, color=RED_ACCENT, bold=True)
    
    # Slide 2: Liquidity Review
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Liquidity Review")
    bullets(slide, [
        "BSL (Buy-Side Liquidity): Above highs - stop losses of shorts",
        "SSL (Sell-Side Liquidity): Below lows - stop losses of longs",
        "Institutions need liquidity to fill large orders",
        "They 'hunt' retail stops before real moves",
        "This session: Advanced liquidity concepts"
    ], 1.2, size=22)
    footer(slide, "Liquidity is the fuel of markets")
    
    # Slide 3: Internal vs External Liquidity
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Internal vs External Liquidity")
    text(slide, "INTERNAL LIQUIDITY", 1.3, 0.5, 4.5, size=22, color=BLUE, bold=True)
    text(slide, "FVGs, Order Blocks, Imbalances", 1.7, 0.5, 4.5, size=16, color=GRAY)
    text(slide, "Inside the current range", 2.0, 0.5, 4.5, size=16, color=GRAY)
    text(slide, "Where price retraces TO", 2.3, 0.5, 4.5, size=16, color=GRAY)
    
    text(slide, "EXTERNAL LIQUIDITY", 1.3, 5, 4.5, size=22, color=GREEN, bold=True)
    text(slide, "Highs, Lows, Equal H/L", 1.7, 5, 4.5, size=16, color=GRAY)
    text(slide, "Outside the current range", 2.0, 5, 4.5, size=16, color=GRAY)
    text(slide, "Where price reaches FOR", 2.3, 5, 4.5, size=16, color=GRAY)
    
    text(slide, "Price moves from internal to external liquidity", 3.5, 0.5, 9, size=20, color=WHITE, bold=True)
    
    # Slide 4: Liquidity Pools
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Identifying Liquidity Pools")
    bullets(slide, [
        "Equal Highs/Lows: Multiple touches = More stops",
        "Trendline touches: Stops placed along trendlines",
        "Previous Day/Week High and Low",
        "Session Highs and Lows (Asian, London, NY)",
        "Psychological round numbers",
        "Old untested highs/lows (virgin liquidity)"
    ], 1.2, size=22)
    footer(slide, "More touches = Bigger pool")
    
    # Slide 5: Liquidity Grabs vs Sweeps
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Liquidity Grab vs Sweep")
    text(slide, "LIQUIDITY GRAB", 1.3, 0.5, 4.5, size=22, color=GREEN, bold=True)
    text(slide, "Quick wick through level", 1.7, 0.5, 4.5, size=16, color=GRAY)
    text(slide, "Immediate reversal", 2.0, 0.5, 4.5, size=16, color=GRAY)
    text(slide, "Clean rejection", 2.3, 0.5, 4.5, size=16, color=GRAY)
    
    text(slide, "LIQUIDITY SWEEP", 1.3, 5, 4.5, size=22, color=RED_ACCENT, bold=True)
    text(slide, "Candle closes beyond level", 1.7, 5, 4.5, size=16, color=GRAY)
    text(slide, "May consolidate before reversing", 2.0, 5, 4.5, size=16, color=GRAY)
    text(slide, "Deeper penetration", 2.3, 5, 4.5, size=16, color=GRAY)
    
    text(slide, "Both are valid - watch for the reversal signal", 3.5, 0.5, 9, size=18, color=WHITE)
    
    # Slide 6: Inducement
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Inducement (The Trap)")
    bullets(slide, [
        "Inducement = Minor liquidity used to trap traders",
        "Creates false breakout/breakdown",
        "Retail sees 'breakout' and enters",
        "Smart money fades the move",
        "The 'obvious' move is usually the trap",
        "Look for inducement BEFORE key levels"
    ], 1.2, size=22)
    footer(slide, "If it looks too easy, it's a trap")
    
    # Slide 7: Liquidity Void
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Liquidity Voids")
    bullets(slide, [
        "Liquidity Void = Area price moved through quickly",
        "Similar to FVG but can span multiple candles",
        "Shows where there's 'no resistance'",
        "Price often returns to fill these voids",
        "Mark large impulsive moves as potential voids",
        "Trade the retracement back into the void"
    ], 1.2, size=22)
    footer(slide, "Voids get filled")
    
    # Slide 8: Trading the Sweep
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Trading Liquidity Sweeps")
    bullets(slide, [
        "1. Identify key liquidity pool (equal highs/lows)",
        "2. Wait for price to sweep the level",
        "3. Look for immediate reversal signal",
        "4. Confirm with BOS or strong rejection candle",
        "5. Enter opposite direction of sweep",
        "6. Target opposite liquidity pool"
    ], 1.2, size=22)
    footer(slide, "Sweep + Reversal = Entry")
    
    # Slide 9: Liquidity Mapping
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Daily Liquidity Mapping")
    bullets(slide, [
        "Mark PDH/PDL (Previous Day High/Low)",
        "Mark PWH/PWL (Previous Week High/Low)",
        "Mark Asian session H/L",
        "Identify equal highs/lows",
        "Note any untested old highs/lows",
        "Update as new levels form"
    ], 1.2, size=22)
    footer(slide, "Map before you trade")
    
    # Slide 10: Key Takeaways
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Key Takeaways")
    bullets(slide, [
        "Internal liquidity = FVGs, OBs (retracement zones)",
        "External liquidity = Highs/Lows (targets)",
        "Equal H/L = Strongest liquidity pools",
        "Inducement traps retail before real move",
        "Sweep + Reversal = High probability entry",
        "Always know where liquidity sits before trading"
    ], 1.2, size=22)
    text(slide, "Follow the liquidity, find the money", 5.5, 0.5, 9, size=20, color=RED_ACCENT, bold=True)
    
    prs.save("/root/clawd/dojo-gains/course-content/M05_Liquidity_DeepDive.pptx")
    print("Saved: M05_Liquidity_DeepDive.pptx")

# ============================================
# 6. FIBONACCI
# ============================================
def create_fibonacci_ppt():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    text(slide, "FIBONACCI", 2.2, 0.5, 9, size=52, color=WHITE, bold=True)
    text(slide, "The Golden Ratio in Trading", 3.1, 0.5, 9, size=26, color=GRAY)
    hline(slide, 3.9, 3.5, 3)
    text(slide, "DOJO GAINS - Mentorship", 4.4, 0.5, 9, size=20, color=RED_ACCENT, bold=True)
    
    # Slide 2: What is Fibonacci?
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "What is Fibonacci?")
    bullets(slide, [
        "Mathematical sequence found throughout nature",
        "0, 1, 1, 2, 3, 5, 8, 13, 21, 34, 55...",
        "Each number is the sum of the two before it",
        "Ratios between numbers approach the 'Golden Ratio'",
        "These ratios appear in market movements",
        "Used to predict support, resistance, and targets"
    ], 1.2, size=22)
    footer(slide, "Math meets market psychology")
    
    # Slide 3: Key Fibonacci Levels
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Key Fibonacci Levels")
    text(slide, "23.6%", 1.3, 0.5, 2.5, size=22, color=GRAY, bold=True)
    text(slide, "Shallow retracement", 1.7, 0.5, 2.5, size=14, color=GRAY)
    
    text(slide, "38.2%", 1.3, 2.5, 2.5, size=22, color=BLUE, bold=True)
    text(slide, "Common retracement", 1.7, 2.5, 2.5, size=14, color=GRAY)
    
    text(slide, "50.0%", 1.3, 5, 2.5, size=22, color=WHITE, bold=True)
    text(slide, "Equilibrium", 1.7, 5, 2.5, size=14, color=GRAY)
    
    text(slide, "61.8%", 1.3, 7.5, 2.5, size=22, color=GREEN, bold=True)
    text(slide, "GOLDEN RATIO", 1.7, 7.5, 2.5, size=14, color=GREEN)
    
    text(slide, "78.6%", 2.4, 0.5, 2.5, size=22, color=RED_ACCENT, bold=True)
    text(slide, "Deep retracement", 2.8, 0.5, 2.5, size=14, color=GRAY)
    
    bullets(slide, [
        "61.8% (Golden Ratio) is the most important level",
        "OTE (Optimal Trade Entry) = 62-79% zone"
    ], 3.8, size=20)
    
    # Slide 4: Fibonacci Retracement
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Fibonacci Retracement")
    bullets(slide, [
        "Used to find potential support/resistance during pullbacks",
        "Draw from swing LOW to swing HIGH (uptrend)",
        "Draw from swing HIGH to swing LOW (downtrend)",
        "Watch for price reaction at Fib levels",
        "Best levels: 50%, 61.8%, 78.6%",
        "OTE zone (62-79%) is highest probability"
    ], 1.2, size=22)
    footer(slide, "Retracements show where to enter")
    
    # Slide 5: Fibonacci Extension
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Fibonacci Extension")
    bullets(slide, [
        "Used to find profit targets beyond the move",
        "Draw from swing LOW to HIGH to pullback LOW (uptrend)",
        "Key extension levels: 127.2%, 161.8%, 200%, 261.8%",
        "161.8% is the most watched extension",
        "Use extensions to set take profit levels",
        "Combine with key S/R for stronger targets"
    ], 1.2, size=22)
    footer(slide, "Extensions show where price may go")
    
    # Slide 6: The OTE Zone
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "OTE - Optimal Trade Entry")
    bullets(slide, [
        "OTE = 62% to 79% Fibonacci retracement zone",
        "This is where smart money typically enters",
        "Deepest retracement before trend continues",
        "Also called 'discount zone' in ICT",
        "Wait for price to enter OTE before looking for entry",
        "Combine with FVG or OB in OTE for A+ setup"
    ], 1.2, size=22)
    text(slide, "OTE + FVG + Killzone = High Probability", 5.5, 0.5, 9, size=20, color=GREEN, bold=True)
    
    # Slide 7: Premium vs Discount
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Premium vs Discount Zones")
    text(slide, "PREMIUM (Above 50%)", 1.3, 0.5, 4.5, size=22, color=RED_ACCENT, bold=True)
    text(slide, "Sell zone in downtrend", 1.7, 0.5, 4.5, size=16, color=GRAY)
    text(slide, "Price is 'expensive'", 2.0, 0.5, 4.5, size=16, color=GRAY)
    text(slide, "Look for shorts here", 2.3, 0.5, 4.5, size=16, color=GRAY)
    
    text(slide, "DISCOUNT (Below 50%)", 1.3, 5, 4.5, size=22, color=GREEN, bold=True)
    text(slide, "Buy zone in uptrend", 1.7, 5, 4.5, size=16, color=GRAY)
    text(slide, "Price is 'cheap'", 2.0, 5, 4.5, size=16, color=GRAY)
    text(slide, "Look for longs here", 2.3, 5, 4.5, size=16, color=GRAY)
    
    text(slide, "Buy at discount, sell at premium", 3.5, 0.5, 9, size=22, color=WHITE, bold=True)
    
    # Slide 8: Using Fib in SMC
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Fibonacci + SMC")
    bullets(slide, [
        "Draw Fib on HTF swing to find retracement zone",
        "Look for FVG or Order Block within OTE",
        "Wait for price to enter OTE during Killzone",
        "Enter at FVG/OB inside OTE",
        "Target: Fib extension (161.8%) or opposite liquidity",
        "This is the institutional entry method"
    ], 1.2, size=22)
    footer(slide, "Fib + FVG + OTE = Sniper entry")
    
    # Slide 9: Common Mistakes
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Fibonacci Mistakes")
    bullets(slide, [
        "Drawing Fib on wrong swing points",
        "Using Fib alone without confluence",
        "Expecting exact reactions at levels",
        "Ignoring the trend direction",
        "Drawing too many Fibs (analysis paralysis)",
        "Not adjusting Fib when structure changes"
    ], 1.2, size=22)
    footer(slide, "Fib is a guide, not a guarantee")
    
    # Slide 10: Key Takeaways
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Key Takeaways")
    bullets(slide, [
        "61.8% is the Golden Ratio - most important level",
        "OTE zone (62-79%) is optimal entry area",
        "Retracements = Entry zones | Extensions = Targets",
        "Premium = Sell zone | Discount = Buy zone",
        "Combine Fib with FVG/OB for best entries",
        "Fib shows WHERE, SMC concepts show WHEN"
    ], 1.2, size=22)
    text(slide, "Fibonacci is your roadmap", 5.5, 0.5, 9, size=22, color=RED_ACCENT, bold=True)
    
    prs.save("/root/clawd/dojo-gains/course-content/M06_Fibonacci.pptx")
    print("Saved: M06_Fibonacci.pptx")

# ============================================
# 7. RISK MANAGEMENT
# ============================================
def create_risk_management_ppt():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    text(slide, "RISK MANAGEMENT", 2.2, 0.5, 9, size=48, color=WHITE, bold=True)
    text(slide, "Protecting Your Capital", 3.1, 0.5, 9, size=26, color=GRAY)
    hline(slide, 3.9, 3.5, 3)
    text(slide, "DOJO GAINS - Mentorship", 4.4, 0.5, 9, size=20, color=RED_ACCENT, bold=True)
    
    # Slide 2: Why Risk Management?
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Why Risk Management Matters")
    bullets(slide, [
        "You can have a winning strategy and still blow your account",
        "One bad trade can wipe out months of gains",
        "Professional traders focus on risk FIRST, profit second",
        "Your job: SURVIVE long enough to succeed",
        "Risk management is not optional - it's everything"
    ], 1.2, size=22)
    footer(slide, "Rule #1: Don't lose money. Rule #2: See Rule #1.")
    
    # Slide 3: The 1% Rule
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "The 1-2% Rule")
    bullets(slide, [
        "NEVER risk more than 1-2% of account on a single trade",
        "$10,000 account = Max $100-200 risk per trade",
        "This allows you to lose 10+ trades and still survive",
        "Drawdowns are inevitable - this keeps them manageable",
        "As account grows, risk amount grows (but % stays same)"
    ], 1.2, size=22)
    text(slide, "1-2% risk = Long-term survival", 5.3, 0.5, 9, size=20, color=GREEN, bold=True)
    
    # Slide 4: Position Sizing
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Position Sizing Formula")
    text(slide, "Position Size = Risk Amount / Stop Distance", 1.4, 0.5, 9, size=24, color=GREEN, bold=True)
    bullets(slide, [
        "Example: $10,000 account, 1% risk = $100",
        "Entry at $50, Stop at $48 = $2 risk per share",
        "Position Size: $100 / $2 = 50 shares",
        "Calculate BEFORE entering every trade",
        "Never adjust size based on 'feeling confident'"
    ], 2.2, size=20)
    footer(slide, "Math > Emotion")
    
    # Slide 5: Risk Reward Ratio
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Risk/Reward Ratio")
    bullets(slide, [
        "R:R = Potential Profit / Potential Loss",
        "1:2 R:R means risking $1 to make $2",
        "Minimum acceptable: 1:1.5 (ideally 1:2 or better)",
        "Higher R:R = Can be profitable with lower win rate"
    ], 1.2, size=22)
    text(slide, "Win Rate Math:", 4.0, 0.5, 9, size=20, color=WHITE, bold=True)
    text(slide, "1:1 R/R needs 50%+ win rate", 4.5, 0.5, 9, size=18, color=GRAY)
    text(slide, "1:2 R/R needs 33%+ win rate", 4.9, 0.5, 9, size=18, color=GRAY)
    text(slide, "1:3 R/R needs 25%+ win rate", 5.3, 0.5, 9, size=18, color=GRAY)
    
    # Slide 6: Stop Loss Placement
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Stop Loss Placement")
    bullets(slide, [
        "Stop should be at a level that INVALIDATES the trade",
        "Below swing low (for longs) / Above swing high (for shorts)",
        "Beyond the FVG or Order Block",
        "Give enough room - avoid getting stopped by noise",
        "NEVER move stop further away",
        "Can move stop to breakeven after 1R profit"
    ], 1.2, size=22)
    footer(slide, "Stop = Where your idea is wrong")
    
    # Slide 7: Daily/Weekly Limits
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Daily & Weekly Limits")
    bullets(slide, [
        "Daily Loss Limit: 2-3% of account",
        "Weekly Loss Limit: 5-6% of account",
        "Hit the limit = DONE trading for that period",
        "No exceptions, no 'making it back'",
        "This prevents emotional trading spirals",
        "Live to trade another day"
    ], 1.2, size=22)
    footer(slide, "Limits save accounts")
    
    # Slide 8: Scaling & Partials
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Taking Partials & Scaling")
    bullets(slide, [
        "Take partial profit at 1R (first target)",
        "Move stop to breakeven after partials",
        "Let remainder run to 2R, 3R targets",
        "This locks in profit while giving room to run",
        "Example: 50% off at 1R, 50% off at 2R",
        "Trailing stop for extended moves"
    ], 1.2, size=22)
    footer(slide, "Bank profits, let winners run")
    
    # Slide 9: Risk Management Rules
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Non-Negotiable Rules")
    bullets(slide, [
        "1. Never risk more than 2% per trade",
        "2. Always use a stop loss",
        "3. Never move stop loss further away",
        "4. Honor daily/weekly loss limits",
        "5. Calculate position size BEFORE entry",
        "6. Minimum 1:2 risk/reward",
        "7. No revenge trading after losses"
    ], 1.2, size=20)
    
    # Slide 10: Key Takeaways
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Key Takeaways")
    bullets(slide, [
        "Risk management > Strategy",
        "1-2% max risk per trade",
        "Position Size = Risk / Stop Distance",
        "Minimum 1:2 R:R on every trade",
        "Daily/weekly limits protect your account",
        "Take partials, move to breakeven, let it run"
    ], 1.2, size=22)
    text(slide, "Protect capital first, profits follow", 5.5, 0.5, 9, size=22, color=RED_ACCENT, bold=True)
    
    prs.save("/root/clawd/dojo-gains/course-content/M07_Risk_Management.pptx")
    print("Saved: M07_Risk_Management.pptx")

# ============================================
# 8. BUILDING A PLAYBOOK
# ============================================
def create_playbook_ppt():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    text(slide, "BUILDING A", 1.8, 0.5, 9, size=28, color=GRAY)
    text(slide, "TRADING PLAYBOOK", 2.5, 0.5, 9, size=44, color=WHITE, bold=True)
    hline(slide, 3.4, 3.5, 3)
    text(slide, "DOJO GAINS - Mentorship", 3.9, 0.5, 9, size=20, color=RED_ACCENT, bold=True)
    text(slide, "Your Personal Trading System", 4.4, 0.5, 9, size=18, color=GRAY)
    
    # Slide 2: What is a Playbook?
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "What is a Trading Playbook?")
    bullets(slide, [
        "A documented set of YOUR proven setups",
        "Clear rules for entry, exit, and management",
        "Based on YOUR backtesting and experience",
        "Removes emotion from trading decisions",
        "Your 'game plan' for every market scenario",
        "Evolves as you learn and grow"
    ], 1.2, size=22)
    footer(slide, "No playbook = No consistency")
    
    # Slide 3: Why You Need One
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Why You Need a Playbook")
    bullets(slide, [
        "Eliminates guessing and FOMO trading",
        "Creates consistency in your approach",
        "Makes journaling and review easier",
        "Identifies what actually works for YOU",
        "Prevents strategy hopping",
        "Professional traders all have one"
    ], 1.2, size=22)
    footer(slide, "Consistency breeds success")
    
    # Slide 4: Playbook Components
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Playbook Components")
    bullets(slide, [
        "1. Setup Name (e.g., 'London Sweep Reversal')",
        "2. Market Conditions (trend, range, volatility)",
        "3. Entry Criteria (specific rules)",
        "4. Exit Criteria (stop loss, targets)",
        "5. Time of Day (when this setup works)",
        "6. Risk Parameters (position size, R:R)",
        "7. Screenshots of Example Trades"
    ], 1.2, size=20)
    
    # Slide 5: Creating Your Setups
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Creating Your Setups")
    bullets(slide, [
        "Start with 2-3 setups MAXIMUM",
        "Each setup should have clear, repeatable rules",
        "Backtest each setup (minimum 50 trades)",
        "Forward test in demo before going live",
        "Document win rate, avg R:R, best times",
        "Only add new setups when current ones mastered"
    ], 1.2, size=22)
    footer(slide, "Master few, not many")
    
    # Slide 6: Example Setup
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Example: London Sweep Setup")
    bullets(slide, [
        "Name: London Sweep Reversal",
        "Time: 2:00 AM - 5:00 AM ET",
        "Condition: Price sweeps Asian H or L",
        "Entry: At FVG or OB after sweep, with BOS",
        "Stop: Beyond sweep wick",
        "Target: Opposite liquidity (Asian L if swept H)",
        "R:R: Minimum 1:2"
    ], 1.2, size=20)
    footer(slide, "This is a complete setup")
    
    # Slide 7: Setup Checklist
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Pre-Trade Checklist")
    bullets(slide, [
        "Is this a setup from my playbook?",
        "Am I trading during the right time?",
        "Is the higher timeframe trend aligned?",
        "Do I have clear entry and stop levels?",
        "Is R:R at least 1:2?",
        "Is my position size correct (1-2% risk)?",
        "Am I emotionally ready to take this trade?"
    ], 1.2, size=20)
    footer(slide, "All YES = Take the trade")
    
    # Slide 8: Tracking & Review
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Tracking & Review Process")
    bullets(slide, [
        "Journal EVERY trade (win or loss)",
        "Screenshot entry, management, exit",
        "Note: What setup? Did you follow rules?",
        "Weekly review: What worked? What didn't?",
        "Monthly: Calculate stats per setup",
        "Quarterly: Update playbook based on data"
    ], 1.2, size=22)
    footer(slide, "Data drives improvement")
    
    # Slide 9: Building Your Playbook
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Building Your Playbook: Steps")
    bullets(slide, [
        "1. Choose 2-3 concepts you understand well",
        "2. Define clear rules for each setup",
        "3. Backtest on historical charts (50+ trades)",
        "4. Document everything in your playbook",
        "5. Forward test in demo for 2-4 weeks",
        "6. Go live with small size, track results",
        "7. Review monthly, refine, repeat"
    ], 1.2, size=20)
    
    # Slide 10: Key Takeaways
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Key Takeaways")
    bullets(slide, [
        "A playbook = Your documented trading system",
        "Start with 2-3 setups, master before adding",
        "Every setup needs clear rules",
        "Use a checklist before every trade",
        "Journal and review religiously",
        "Your playbook evolves with you"
    ], 1.2, size=22)
    text(slide, "Your playbook is your edge", 5.5, 0.5, 9, size=22, color=RED_ACCENT, bold=True)
    
    prs.save("/root/clawd/dojo-gains/course-content/M08_Building_Playbook.pptx")
    print("Saved: M08_Building_Playbook.pptx")

# ============================================
# 9. INTRO TO CHARTING
# ============================================
def create_charting_intro_ppt():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    text(slide, "INTRO TO CHARTING", 2.2, 0.5, 9, size=44, color=WHITE, bold=True)
    text(slide, "Reading the Market's Language", 3.1, 0.5, 9, size=26, color=GRAY)
    hline(slide, 3.9, 3.5, 3)
    text(slide, "DOJO GAINS - Mentorship", 4.4, 0.5, 9, size=20, color=RED_ACCENT, bold=True)
    
    # Slide 2: What is a Chart?
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "What is a Chart?")
    bullets(slide, [
        "Visual representation of price over time",
        "Shows the battle between buyers and sellers",
        "Each candle/bar = price action for a time period",
        "Charts reveal patterns that repeat",
        "The foundation of technical analysis",
        "Your window into market psychology"
    ], 1.2, size=22)
    footer(slide, "Charts are the trader's map")
    
    # Slide 3: Chart Types
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Chart Types")
    text(slide, "CANDLESTICK", 1.3, 0.5, 4.5, size=22, color=GREEN, bold=True)
    text(slide, "Most popular, shows OHLC", 1.7, 0.5, 4.5, size=16, color=GRAY)
    text(slide, "Best for price action trading", 2.0, 0.5, 4.5, size=16, color=GRAY)
    
    text(slide, "BAR CHART", 1.3, 5, 4.5, size=22, color=BLUE, bold=True)
    text(slide, "Similar to candle, different visual", 1.7, 5, 4.5, size=16, color=GRAY)
    text(slide, "Used by some futures traders", 2.0, 5, 4.5, size=16, color=GRAY)
    
    text(slide, "LINE CHART", 2.7, 0.5, 4.5, size=22, color=GRAY, bold=True)
    text(slide, "Only shows close prices", 3.1, 0.5, 4.5, size=16, color=GRAY)
    text(slide, "Good for seeing trends simply", 3.4, 0.5, 4.5, size=16, color=GRAY)
    
    text(slide, "HEIKIN ASHI", 2.7, 5, 4.5, size=22, color=YELLOW, bold=True)
    text(slide, "Smoothed candles", 3.1, 5, 4.5, size=16, color=GRAY)
    text(slide, "Good for trend identification", 3.4, 5, 4.5, size=16, color=GRAY)
    
    footer(slide, "We use candlestick charts")
    
    # Slide 4: Timeframes
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Understanding Timeframes")
    bullets(slide, [
        "Each candle = one period of time",
        "1 minute chart: 1 candle = 1 minute",
        "Daily chart: 1 candle = 1 full day",
        "Higher TF = less noise, more reliable",
        "Lower TF = more detail, more noise"
    ], 1.2, size=22)
    text(slide, "Day Trading: 1m, 5m, 15m", 4.5, 0.5, 9, size=18, color=GREEN)
    text(slide, "Swing Trading: 1H, 4H, Daily", 4.9, 0.5, 9, size=18, color=BLUE)
    text(slide, "Position Trading: Daily, Weekly", 5.3, 0.5, 9, size=18, color=YELLOW)
    
    # Slide 5: Multi-Timeframe Analysis
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Multi-Timeframe Analysis")
    bullets(slide, [
        "Always check higher timeframe first",
        "HTF = Direction and key levels",
        "LTF = Entry timing and precision",
        "Example: Daily for trend, 15m for entry",
        "Don't fight the higher timeframe trend",
        "Alignment across timeframes = higher probability"
    ], 1.2, size=22)
    footer(slide, "Top-down analysis is key")
    
    # Slide 6: Essential Chart Elements
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Essential Chart Elements")
    bullets(slide, [
        "PRICE: The candles/bars (most important)",
        "VOLUME: Confirms price moves",
        "SUPPORT/RESISTANCE: Key levels",
        "TREND: Direction of the market",
        "PATTERNS: Repeating formations",
        "Keep charts clean - less is more"
    ], 1.2, size=22)
    footer(slide, "Price is king, everything else supports it")
    
    # Slide 7: Setting Up Your Charts
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Setting Up Your Charts")
    bullets(slide, [
        "Use dark background (easier on eyes)",
        "Green candles = bullish, Red = bearish",
        "Mark key levels (S/R, highs, lows)",
        "Add 1-2 key indicators max",
        "Save templates for different strategies",
        "Keep it simple and clean"
    ], 1.2, size=22)
    footer(slide, "Clean charts = Clear thinking")
    
    # Slide 8: Key Takeaways
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Key Takeaways")
    bullets(slide, [
        "Charts show price over time",
        "Candlestick charts are standard",
        "Higher timeframes = more reliable",
        "Always use multi-timeframe analysis",
        "Keep charts clean and simple",
        "Price is the primary focus"
    ], 1.2, size=22)
    text(slide, "The chart tells you everything you need to know", 5.5, 0.5, 9, size=18, color=RED_ACCENT)
    
    prs.save("/root/clawd/dojo-gains/course-content/M09_Intro_To_Charting.pptx")
    print("Saved: M09_Intro_To_Charting.pptx")

# ============================================
# MAIN
# ============================================
if __name__ == "__main__":
    create_mindset_ppt()
    create_fvg_deep_dive_ppt()
    create_divergences_ppt()
    create_candle_deep_dive_ppt()
    create_liquidity_deep_dive_ppt()
    create_fibonacci_ppt()
    create_risk_management_ppt()
    create_playbook_ppt()
    create_charting_intro_ppt()
    print("\nAll 9 mentorship presentations created!")
