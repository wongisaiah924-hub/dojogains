#!/usr/bin/env python3
"""
Dojo Gains - SMC/ICT Core Pillars Presentations
Creates: Liquidity, Time, Imbalances, Intro to Strategy
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
DARK_BG = RGBColor(15, 15, 26)
RED_ACCENT = RGBColor(230, 57, 70)
GREEN = RGBColor(34, 197, 94)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(156, 163, 175)
BLUE = RGBColor(59, 130, 246)

SLIDE_W = 10
SLIDE_H = 7.5

def set_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK_BG

def text(slide, txt, top, left, width, size=18, color=WHITE, bold=False, align="center"):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = PP_ALIGN.CENTER if align == "center" else PP_ALIGN.LEFT
    return box

def title(slide, txt, top=0.4):
    text(slide, txt, top, 0.5, 9, size=40, color=WHITE, bold=True)

def bullets(slide, points, top, size=20):
    box = slide.shapes.add_textbox(Inches(0.75), Inches(top), Inches(8.5), Inches(5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, pt in enumerate(points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•   {pt}"
        p.font.size = Pt(size)
        p.font.color.rgb = WHITE
        p.space_before = Pt(6)
        p.space_after = Pt(10)

def footer(slide, txt):
    text(slide, txt, 6.9, 0.5, 9, size=14, color=GRAY)

def hline(slide, top, left, width, color=RED_ACCENT):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.02))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()

# ============================================
# PRESENTATION 1: LIQUIDITY
# ============================================
def create_liquidity_ppt():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    text(slide, "LIQUIDITY", 2.2, 0.5, 9, size=52, color=WHITE, bold=True)
    text(slide, "Where Smart Money Finds Its Fuel", 3.1, 0.5, 9, size=26, color=GRAY)
    hline(slide, 3.9, 3.5, 3)
    text(slide, "DOJO GAINS - Core Pillar 1", 4.4, 0.5, 9, size=20, color=RED_ACCENT, bold=True)
    
    # Slide 2: What is Liquidity?
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "What is Liquidity?")
    bullets(slide, [
        "Liquidity = Orders resting in the market (stop losses & pending orders)",
        "Institutions need liquidity to fill large positions",
        "Retail traders place stops at predictable levels",
        "Smart money 'hunts' these stop losses to get their orders filled",
        "Understanding liquidity = understanding where price will go"
    ], 1.2, size=22)
    footer(slide, "Liquidity is the fuel that moves markets")
    
    # Slide 3: Buy-Side Liquidity (BSL)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Buy-Side Liquidity (BSL)")
    bullets(slide, [
        "Located ABOVE swing highs, equal highs, and resistance levels",
        "Consists of buy stop orders from traders who are SHORT",
        "When triggered, buy stops become market BUY orders",
        "Institutions push price UP to grab this liquidity",
        "Often leads to a REVERSAL after the sweep"
    ], 1.2, size=22)
    
    # Visual representation
    text(slide, "RESISTANCE / SWING HIGH", 4.8, 0.5, 9, size=16, color=GRAY)
    hline(slide, 5.15, 2, 6, color=GRAY)
    text(slide, "Buy Stops Resting Here (BSL)", 5.3, 0.5, 9, size=18, color=GREEN, bold=True)
    text(slide, "Price sweeps above, grabs liquidity, then reverses", 5.8, 0.5, 9, size=14, color=GRAY)
    
    # Slide 4: Sell-Side Liquidity (SSL)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Sell-Side Liquidity (SSL)")
    bullets(slide, [
        "Located BELOW swing lows, equal lows, and support levels",
        "Consists of sell stop orders from traders who are LONG",
        "When triggered, sell stops become market SELL orders",
        "Institutions push price DOWN to grab this liquidity",
        "Often leads to a REVERSAL after the sweep"
    ], 1.2, size=22)
    
    # Visual representation
    text(slide, "Sell Stops Resting Here (SSL)", 4.8, 0.5, 9, size=18, color=RED_ACCENT, bold=True)
    hline(slide, 5.25, 2, 6, color=GRAY)
    text(slide, "SUPPORT / SWING LOW", 5.4, 0.5, 9, size=16, color=GRAY)
    text(slide, "Price sweeps below, grabs liquidity, then reverses", 5.9, 0.5, 9, size=14, color=GRAY)
    
    # Slide 5: Where Liquidity Pools Form
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Where Liquidity Pools Form")
    bullets(slide, [
        "Previous Day High (PDH) / Previous Day Low (PDL)",
        "Previous Week High (PWH) / Previous Week Low (PWL)",
        "Swing Highs and Swing Lows",
        "Equal Highs and Equal Lows (double tops/bottoms)",
        "Round psychological numbers ($100, $50, etc.)",
        "Session highs and lows (Asian, London, NY)"
    ], 1.2, size=22)
    footer(slide, "More touches = More stops = Bigger liquidity pool")
    
    # Slide 6: Liquidity Sweeps
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Liquidity Sweeps (Stop Hunts)")
    bullets(slide, [
        "A liquidity sweep happens when price takes out a key level",
        "Price breaks the high/low, triggers stops, then reverses",
        "The 'wick' often shows the sweep on the candle",
        "Look for rejection after the sweep (long wicks)",
        "This is smart money getting filled before the real move"
    ], 1.2, size=22)
    
    text(slide, "Sweep = Trap = Reversal Incoming", 5.5, 0.5, 9, size=22, color=RED_ACCENT, bold=True)
    
    # Slide 7: Trading Liquidity
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "How to Trade Liquidity")
    bullets(slide, [
        "1. Identify where liquidity is resting (highs/lows)",
        "2. Wait for price to sweep that liquidity",
        "3. Look for confirmation (rejection candle, BOS)",
        "4. Enter in the opposite direction of the sweep",
        "5. Target the opposite liquidity pool"
    ], 1.2, size=22)
    footer(slide, "Don't trade the breakout - trade the reversal after the sweep")
    
    # Slide 8: Liquidity Example
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Liquidity Sweep Example")
    bullets(slide, [
        "Price creates a swing high - retail goes short with stops above",
        "Market rallies and takes out the high (liquidity grab)",
        "Price immediately reverses with strong selling",
        "Smart money sold into the buy stops getting triggered",
        "Result: Price drops and targets sell-side liquidity below"
    ], 1.2, size=22)
    footer(slide, "The hunt is real - be the hunter, not the hunted")
    
    # Slide 9: Key Takeaways
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Key Takeaways")
    bullets(slide, [
        "BSL (Buy-Side) = Above highs = Target when bearish",
        "SSL (Sell-Side) = Below lows = Target when bullish",
        "Smart money sweeps liquidity before making real moves",
        "Equal highs/lows are liquidity magnets",
        "Don't place stops at obvious levels",
        "Trade WITH smart money, not against them"
    ], 1.2, size=22)
    
    text(slide, "Liquidity is the lifeblood of institutional trading", 5.8, 0.5, 9, size=18, color=RED_ACCENT)
    
    prs.save("/root/clawd/dojo-gains/course-content/01_Liquidity_DojoGains.pptx")
    print("Saved: 01_Liquidity_DojoGains.pptx")

# ============================================
# PRESENTATION 2: TIME
# ============================================
def create_time_ppt():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    text(slide, "TIME", 2.2, 0.5, 9, size=52, color=WHITE, bold=True)
    text(slide, "When Smart Money Moves", 3.1, 0.5, 9, size=26, color=GRAY)
    hline(slide, 3.9, 3.5, 3)
    text(slide, "DOJO GAINS - Core Pillar 2", 4.4, 0.5, 9, size=20, color=RED_ACCENT, bold=True)
    
    # Slide 2: Why Time Matters
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Why Time Matters")
    bullets(slide, [
        "Markets are not active 24/7 - there are windows of opportunity",
        "Institutions trade during specific hours (Killzones)",
        "Volume and volatility peak during certain times",
        "Wrong time = low probability setups, fakeouts, chop",
        "Right time = clean moves, follow-through, real setups"
    ], 1.2, size=22)
    footer(slide, "Time is the invisible edge most traders ignore")
    
    # Slide 3: The Four Killzones
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "The Four Killzones")
    
    text(slide, "ASIAN SESSION", 1.3, 0.5, 4.5, size=20, color=BLUE, bold=True)
    text(slide, "7:00 PM - 12:00 AM ET", 1.7, 0.5, 4.5, size=16, color=GRAY)
    text(slide, "Range building, liquidity setup", 2.05, 0.5, 4.5, size=14, color=GRAY)
    
    text(slide, "LONDON SESSION", 2.7, 0.5, 4.5, size=20, color=GREEN, bold=True)
    text(slide, "2:00 AM - 5:00 AM ET", 3.1, 0.5, 4.5, size=16, color=GRAY)
    text(slide, "First major move, trend initiation", 3.45, 2.7, 4.5, size=14, color=GRAY)
    
    text(slide, "NY AM SESSION", 1.3, 5, 4.5, size=20, color=RED_ACCENT, bold=True)
    text(slide, "8:30 AM - 11:00 AM ET", 1.7, 5, 4.5, size=16, color=GRAY)
    text(slide, "Highest volume, strongest moves", 2.05, 5, 4.5, size=14, color=GRAY)
    
    text(slide, "NY PM SESSION", 2.7, 5, 4.5, size=20, color=WHITE, bold=True)
    text(slide, "1:00 PM - 4:00 PM ET", 3.1, 5, 4.5, size=16, color=GRAY)
    text(slide, "Continuation or reversal", 3.45, 5, 4.5, size=14, color=GRAY)
    
    footer(slide, "All times in New York (Eastern) Time")
    
    # Slide 4: Asian Session
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Asian Session (7 PM - 12 AM ET)")
    bullets(slide, [
        "Generally the LOWEST volatility session",
        "Price often consolidates and builds a range",
        "Creates liquidity above and below the range",
        "Sets up the 'Asian High' and 'Asian Low'",
        "Smart money accumulates positions quietly",
        "Best use: Mark the range, prepare for London"
    ], 1.2, size=22)
    footer(slide, "The calm before the storm")
    
    # Slide 5: London Session
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "London Killzone (2 AM - 5 AM ET)")
    bullets(slide, [
        "First major liquidity injection of the day",
        "Often sweeps Asian session high OR low",
        "Sets the direction for the day (London = True Move)",
        "High probability setups form here",
        "Best for: EUR, GBP pairs, indices",
        "Watch for manipulation then reversal"
    ], 1.2, size=22)
    footer(slide, "London sets the tone - pay attention")
    
    # Slide 6: NY AM Session
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "New York AM (8:30 AM - 11 AM ET)")
    bullets(slide, [
        "HIGHEST volume session - most opportunities",
        "US economic data releases (8:30 AM)",
        "Often continues London's move OR reverses it",
        "Watch for NY Open manipulation (9:30 AM)",
        "Best session for US stocks, indices, USD pairs",
        "Power Hour: 9:30 - 10:30 AM"
    ], 1.2, size=22)
    footer(slide, "This is where the money is made")
    
    # Slide 7: NY PM Session
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "New York PM (1 PM - 4 PM ET)")
    bullets(slide, [
        "Volume decreases after lunch (11 AM - 1 PM is dead)",
        "Often sees reversals or trend continuation",
        "Last hour (3-4 PM) can have strong moves",
        "Institutions closing positions",
        "Good for: Swing trade entries, position adjustments",
        "Avoid: 11 AM - 1 PM (lunch = chop)"
    ], 1.2, size=22)
    footer(slide, "Finish strong or step aside")
    
    # Slide 8: Time-Based Strategy
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Time-Based Trading Strategy")
    bullets(slide, [
        "1. Mark the Asian range (high and low)",
        "2. At London open, watch for sweep of Asian H/L",
        "3. Enter after sweep + confirmation",
        "4. If no setup in London, wait for NY AM",
        "5. Avoid trading outside Killzones",
        "6. Close or manage by NY lunch"
    ], 1.2, size=22)
    footer(slide, "Patience during off-hours = profits during Killzones")
    
    # Slide 9: Key Takeaways
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Key Takeaways")
    bullets(slide, [
        "Asian = Range building (mark H/L)",
        "London = First real move (2-5 AM ET)",
        "NY AM = Highest volume (8:30-11 AM ET)",
        "NY PM = Continuation/reversal (1-4 PM ET)",
        "Avoid: Asian (unless ranging) and NY Lunch",
        "Best setups happen IN the Killzones"
    ], 1.2, size=22)
    
    text(slide, "Trade the right time, not all the time", 5.8, 0.5, 9, size=18, color=RED_ACCENT)
    
    prs.save("/root/clawd/dojo-gains/course-content/02_Time_DojoGains.pptx")
    print("Saved: 02_Time_DojoGains.pptx")

# ============================================
# PRESENTATION 3: IMBALANCES
# ============================================
def create_imbalances_ppt():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    text(slide, "IMBALANCES", 2.2, 0.5, 9, size=52, color=WHITE, bold=True)
    text(slide, "Fair Value Gaps & Market Inefficiency", 3.1, 0.5, 9, size=26, color=GRAY)
    hline(slide, 3.9, 3.5, 3)
    text(slide, "DOJO GAINS - Core Pillar 3", 4.4, 0.5, 9, size=20, color=RED_ACCENT, bold=True)
    
    # Slide 2: What is an Imbalance?
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "What is an Imbalance?")
    bullets(slide, [
        "An imbalance occurs when price moves too fast",
        "Creates a gap where not all orders were filled",
        "Also called: Fair Value Gap (FVG), Inefficiency",
        "Shows aggressive buying OR selling",
        "Price often returns to 'fill' or 'rebalance' the gap",
        "Key concept in ICT and Smart Money trading"
    ], 1.2, size=22)
    footer(slide, "Markets seek balance - imbalances get filled")
    
    # Slide 3: Identifying FVG
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Identifying a Fair Value Gap")
    bullets(slide, [
        "Look at 3 consecutive candles",
        "Middle candle = the 'displacement' candle (large body)",
        "Gap = space between Candle 1's wick and Candle 3's wick",
        "BULLISH FVG: Candle 3's LOW > Candle 1's HIGH",
        "BEARISH FVG: Candle 3's HIGH < Candle 1's LOW",
        "The gap area is your FVG zone"
    ], 1.2, size=22)
    footer(slide, "3 candles, 1 gap = FVG")
    
    # Slide 4: Bullish FVG
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Bullish Fair Value Gap")
    bullets(slide, [
        "Forms during strong upward moves",
        "Gap between Candle 1 HIGH and Candle 3 LOW",
        "Shows aggressive BUYING pressure",
        "Price often retraces DOWN to fill this gap",
        "FVG becomes a SUPPORT zone",
        "Look for LONG entries when price returns to FVG"
    ], 1.2, size=22)
    
    text(slide, "Bullish FVG = Buy Zone on Pullback", 5.5, 0.5, 9, size=20, color=GREEN, bold=True)
    
    # Slide 5: Bearish FVG
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Bearish Fair Value Gap")
    bullets(slide, [
        "Forms during strong downward moves",
        "Gap between Candle 1 LOW and Candle 3 HIGH",
        "Shows aggressive SELLING pressure",
        "Price often retraces UP to fill this gap",
        "FVG becomes a RESISTANCE zone",
        "Look for SHORT entries when price returns to FVG"
    ], 1.2, size=22)
    
    text(slide, "Bearish FVG = Sell Zone on Rally", 5.5, 0.5, 9, size=20, color=RED_ACCENT, bold=True)
    
    # Slide 6: FVG Fill Levels
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "FVG Fill Levels")
    bullets(slide, [
        "50% (Equilibrium/CE): Most common fill level",
        "100% (Full Fill): Price fills entire gap",
        "Partial Fill: Price touches FVG but doesn't fill 50%",
        "The stronger the trend, the less likely a full fill",
        "Use 50% as your primary entry zone",
        "Mark the FVG on your chart and wait for price to return"
    ], 1.2, size=22)
    footer(slide, "50% is the sweet spot for entries")
    
    # Slide 7: Trading FVGs
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "How to Trade Imbalances")
    bullets(slide, [
        "1. Identify the trend direction",
        "2. Look for displacement (big candle creating FVG)",
        "3. Mark the FVG zone on your chart",
        "4. Wait for price to retrace to the FVG",
        "5. Enter at the 50% level with confirmation",
        "6. Stop below/above the FVG, target next liquidity"
    ], 1.2, size=22)
    footer(slide, "FVG + Trend alignment = High probability")
    
    # Slide 8: FVG Confluence
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "FVG + Confluence = A+ Setup")
    bullets(slide, [
        "FVG at an Order Block = strong",
        "FVG at a key support/resistance = strong",
        "FVG in a Killzone time = stronger",
        "FVG after liquidity sweep = strongest",
        "Multiple FVGs overlapping = mega zone",
        "Always combine FVG with other SMC concepts"
    ], 1.2, size=22)
    footer(slide, "More confluence = higher probability")
    
    # Slide 9: Key Takeaways
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Key Takeaways")
    bullets(slide, [
        "FVG = Gap created by fast price movement",
        "3-candle formation: Check Candle 1 & 3 wicks",
        "Bullish FVG = Support zone for longs",
        "Bearish FVG = Resistance zone for shorts",
        "Target the 50% (equilibrium) for entries",
        "Price seeks to fill imbalances before continuing"
    ], 1.2, size=22)
    
    text(slide, "Imbalances are magnets - price comes back", 5.8, 0.5, 9, size=18, color=RED_ACCENT)
    
    prs.save("/root/clawd/dojo-gains/course-content/03_Imbalances_DojoGains.pptx")
    print("Saved: 03_Imbalances_DojoGains.pptx")

# ============================================
# PRESENTATION 4: INTRO TO STRATEGY
# ============================================
def create_strategy_ppt():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    text(slide, "INTRODUCTION TO", 1.8, 0.5, 9, size=28, color=GRAY)
    text(slide, "MY TRADING STRATEGY", 2.5, 0.5, 9, size=44, color=WHITE, bold=True)
    hline(slide, 3.5, 3.5, 3)
    text(slide, "DOJO GAINS", 4.0, 0.5, 9, size=20, color=RED_ACCENT, bold=True)
    text(slide, "Putting the Core Pillars Together", 4.5, 0.5, 9, size=18, color=GRAY)
    
    # Slide 2: The Three Pillars
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "The Three Core Pillars")
    
    text(slide, "1. LIQUIDITY", 1.5, 0.5, 9, size=24, color=GREEN, bold=True)
    text(slide, "Where is the money? Where are stops resting?", 1.95, 0.5, 9, size=18, color=GRAY)
    
    text(slide, "2. TIME", 2.6, 0.5, 9, size=24, color=BLUE, bold=True)
    text(slide, "When do institutions move? Trade the Killzones.", 3.05, 0.5, 9, size=18, color=GRAY)
    
    text(slide, "3. IMBALANCES", 3.7, 0.5, 9, size=24, color=RED_ACCENT, bold=True)
    text(slide, "Where did price move too fast? FVGs are entry zones.", 4.15, 0.5, 9, size=18, color=GRAY)
    
    text(slide, "Combine all three = High probability trades", 5.3, 0.5, 9, size=22, color=WHITE, bold=True)
    
    # Slide 3: The Strategy Overview
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Strategy Overview")
    bullets(slide, [
        "Wait for a KILLZONE (London or NY AM)",
        "Identify where LIQUIDITY is resting (BSL/SSL)",
        "Watch for a SWEEP of that liquidity",
        "Look for IMBALANCE (FVG) after the sweep",
        "Enter at the FVG with defined risk",
        "Target the opposite liquidity pool"
    ], 1.2, size=22)
    footer(slide, "Simple framework, powerful results")
    
    # Slide 4: Step 1 - Mark the Levels
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Step 1: Mark Key Levels")
    bullets(slide, [
        "Draw previous day high (PDH) and low (PDL)",
        "Mark Asian session high and low",
        "Identify swing highs and swing lows",
        "Note equal highs/lows (liquidity magnets)",
        "Mark any unfilled FVGs from previous sessions",
        "These are your targets and entry zones"
    ], 1.2, size=22)
    footer(slide, "Preparation happens BEFORE the Killzone")
    
    # Slide 5: Step 2 - Wait for Time
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Step 2: Wait for Killzone")
    bullets(slide, [
        "Do NOT trade during Asian (unless ranging)",
        "LONDON KILLZONE: 2:00 AM - 5:00 AM ET",
        "NY AM KILLZONE: 8:30 AM - 11:00 AM ET",
        "Watch for manipulation at session opens",
        "Be patient - the setup will come",
        "No setup in Killzone = No trade"
    ], 1.2, size=22)
    footer(slide, "Patience is profit")
    
    # Slide 6: Step 3 - Watch for Sweep
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Step 3: Wait for Liquidity Sweep")
    bullets(slide, [
        "Watch for price to take out a key level",
        "Sweep of Asian high/low is common at London open",
        "Look for wicks above/below the level",
        "The sweep triggers retail stops",
        "Smart money gets filled during the sweep",
        "After the sweep, prepare for reversal"
    ], 1.2, size=22)
    footer(slide, "The sweep is your signal")
    
    # Slide 7: Step 4 - Find the Entry
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Step 4: Enter at the FVG")
    bullets(slide, [
        "After the sweep, look for displacement",
        "Mark the Fair Value Gap that forms",
        "Wait for price to retrace to the FVG",
        "Enter at 50% of the FVG (equilibrium)",
        "Stop loss: Beyond the FVG",
        "Confirmation: Rejection candle at FVG"
    ], 1.2, size=22)
    footer(slide, "FVG = Your entry zone")
    
    # Slide 8: Step 5 - Manage the Trade
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Step 5: Target & Manage")
    bullets(slide, [
        "Target: Opposite liquidity pool (BSL/SSL)",
        "Take partials at 50% and equilibrium points",
        "Move stop to breakeven after first target",
        "Trail stop as price moves in your favor",
        "Exit or manage by NY lunch (11 AM)",
        "Risk:Reward minimum 1:2 (ideally 1:3+)"
    ], 1.2, size=22)
    footer(slide, "Plan your exit before you enter")
    
    # Slide 9: Example Trade
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Example Trade Setup")
    bullets(slide, [
        "Asian session creates a range (mark H/L)",
        "London opens at 2 AM, sweeps Asian LOW",
        "Strong bullish candle creates FVG",
        "Price retraces to FVG - ENTER LONG",
        "Stop below FVG, target Asian HIGH / PDH",
        "Price hits target - take profits"
    ], 1.2, size=22)
    footer(slide, "Liquidity + Time + Imbalance = Trade")
    
    # Slide 10: Key Rules
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Key Trading Rules")
    bullets(slide, [
        "ONLY trade during Killzones",
        "WAIT for liquidity sweep before entering",
        "ENTER at FVG with confirmation",
        "RISK max 1-2% per trade",
        "TARGET the opposite liquidity",
        "NO revenge trading - one good trade per day is enough"
    ], 1.2, size=22)
    
    text(slide, "Discipline over everything", 5.8, 0.5, 9, size=20, color=RED_ACCENT, bold=True)
    
    # Slide 11: Summary
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Strategy Summary")
    
    text(slide, "LIQUIDITY tells you WHERE price will go", 1.4, 0.5, 9, size=22, color=GREEN)
    text(slide, "TIME tells you WHEN to look for setups", 2.1, 0.5, 9, size=22, color=BLUE)
    text(slide, "IMBALANCES tell you WHERE to enter", 2.8, 0.5, 9, size=22, color=RED_ACCENT)
    
    text(slide, "The Strategy:", 3.8, 0.5, 9, size=20, color=WHITE, bold=True)
    text(slide, "Mark levels > Wait for Killzone > Watch for sweep >", 4.3, 0.5, 9, size=18, color=GRAY)
    text(slide, "Enter at FVG > Target opposite liquidity", 4.7, 0.5, 9, size=18, color=GRAY)
    
    text(slide, "dojogains.com", 5.8, 0.5, 9, size=28, color=RED_ACCENT, bold=True)
    text(slide, "No hype. No gurus. Just growth.", 6.4, 0.5, 9, size=16, color=GRAY)
    
    prs.save("/root/clawd/dojo-gains/course-content/04_Strategy_Intro_DojoGains.pptx")
    print("Saved: 04_Strategy_Intro_DojoGains.pptx")

if __name__ == "__main__":
    create_liquidity_ppt()
    create_time_ppt()
    create_imbalances_ppt()
    create_strategy_ppt()
    print("\nAll 4 presentations created!")
