#!/usr/bin/env python3
"""
Dojo Gains - Enhanced Visual Presentations
Adds diagrams, examples, and image placeholders
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
DARK_BG = RGBColor(15, 15, 26)
RED_ACCENT = RGBColor(230, 57, 70)
GREEN = RGBColor(34, 197, 94)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(156, 163, 175)
BLUE = RGBColor(59, 130, 246)
YELLOW = RGBColor(234, 179, 8)
DARK_GRAY = RGBColor(45, 45, 60)

SLIDE_W = 10
SLIDE_H = 7.5

def set_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK_BG

def text(slide, txt, top, left, width, size=18, color=WHITE, bold=False, align="center"):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = PP_ALIGN.CENTER if align == "center" else PP_ALIGN.LEFT
    return box

def title(slide, txt, top=0.4):
    text(slide, txt, top, 0.5, 9, size=40, color=WHITE, bold=True)

def bullets(slide, points, top, size=20):
    box = slide.shapes.add_textbox(Inches(0.75), Inches(top), Inches(8.5), Inches(5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, pt in enumerate(points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•   {pt}"
        p.font.size = Pt(size)
        p.font.color.rgb = WHITE
        p.space_before = Pt(6)
        p.space_after = Pt(10)

def footer(slide, txt):
    text(slide, txt, 6.9, 0.5, 9, size=14, color=GRAY)

def hline(slide, top, left, width, color=RED_ACCENT):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.02))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()

def candle(slide, cx, top, bullish=True, body=1.0, up_wick=0.3, dn_wick=0.3, scale=1.0):
    """Draw a candlestick"""
    color = GREEN if bullish else RED_ACCENT
    bw = 0.4 * scale
    ww = 0.05 * scale
    body *= scale
    up_wick *= scale
    dn_wick *= scale
    
    if up_wick > 0.01:
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
            Inches(cx - ww/2), Inches(top), Inches(ww), Inches(up_wick))
        s.fill.solid()
        s.fill.fore_color.rgb = color
        s.line.fill.background()
    
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(cx - bw/2), Inches(top + up_wick), Inches(bw), Inches(body))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    
    if dn_wick > 0.01:
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(cx - ww/2), Inches(top + up_wick + body), Inches(ww), Inches(dn_wick))
        s.fill.solid()
        s.fill.fore_color.rgb = color
        s.line.fill.background()

def candle_bottom(slide, cx, baseline, bullish=True, body=1.0, up_wick=0.3, dn_wick=0.3, scale=1.0):
    """Draw candlestick with bottom at baseline"""
    color = GREEN if bullish else RED_ACCENT
    bw = 0.4 * scale
    ww = 0.05 * scale
    body *= scale
    up_wick *= scale
    dn_wick *= scale
    total = up_wick + body + dn_wick
    top = baseline - total
    
    if up_wick > 0.01:
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
            Inches(cx - ww/2), Inches(top), Inches(ww), Inches(up_wick))
        s.fill.solid()
        s.fill.fore_color.rgb = color
        s.line.fill.background()
    
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(cx - bw/2), Inches(top + up_wick), Inches(bw), Inches(body))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    
    if dn_wick > 0.01:
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(cx - ww/2), Inches(top + up_wick + body), Inches(ww), Inches(dn_wick))
        s.fill.solid()
        s.fill.fore_color.rgb = color
        s.line.fill.background()

def image_placeholder(slide, top, left, width, height, label="INSERT CHART EXAMPLE"):
    """Add a placeholder box for images"""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
        Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid()
    s.fill.fore_color.rgb = DARK_GRAY
    s.line.color.rgb = GRAY
    s.line.width = Pt(2)
    
    # Add label text
    box = slide.shapes.add_textbox(Inches(left), Inches(top + height/2 - 0.2), Inches(width), Inches(0.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

def arrow_down(slide, cx, top, length, color=WHITE):
    """Draw a down arrow"""
    # Vertical line
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(cx - 0.02), Inches(top), Inches(0.04), Inches(length - 0.1))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    # Arrow head
    s = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
        Inches(cx - 0.1), Inches(top + length - 0.15), Inches(0.2), Inches(0.15))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.rotation = 180

def arrow_up(slide, cx, top, length, color=WHITE):
    """Draw an up arrow"""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(cx - 0.02), Inches(top + 0.1), Inches(0.04), Inches(length - 0.1))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
        Inches(cx - 0.1), Inches(top), Inches(0.2), Inches(0.15))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()


# ============================================
# LIQUIDITY - WITH VISUALS
# ============================================
def create_liquidity_visual():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    text(slide, "LIQUIDITY", 2.2, 0.5, 9, size=52, color=WHITE, bold=True)
    text(slide, "Where Smart Money Finds Its Fuel", 3.1, 0.5, 9, size=26, color=GRAY)
    hline(slide, 3.9, 3.5, 3)
    text(slide, "DOJO GAINS - Core Pillar 1", 4.4, 0.5, 9, size=20, color=RED_ACCENT, bold=True)
    
    # Slide 2: What is Liquidity - with diagram
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "What is Liquidity?")
    
    # Left side - text
    bullets(slide, [
        "Liquidity = Resting orders (stop losses)",
        "Institutions need liquidity to fill positions",
        "Retail places stops at predictable levels",
        "Smart money 'hunts' these stops"
    ], 1.2, size=20)
    
    # Right side - simple diagram
    text(slide, "LIQUIDITY CONCEPT", 1.2, 6, 3.5, size=14, color=GRAY, bold=True)
    
    # Draw swing high with stops above
    hline(slide, 2.5, 6.2, 2.5, color=GRAY)
    text(slide, "Swing High", 2.6, 6.2, 2.5, size=12, color=GRAY)
    
    # Stops above
    for i in range(5):
        s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.4 + i*0.4), Inches(2.1), Inches(0.15), Inches(0.15))
        s.fill.solid()
        s.fill.fore_color.rgb = RED_ACCENT
        s.line.fill.background()
    text(slide, "Stop Losses (BSL)", 1.8, 6.2, 2.5, size=11, color=RED_ACCENT)
    
    # Price line going up to grab them
    arrow_up(slide, 7.5, 3.0, 0.6, GREEN)
    text(slide, "Price sweeps", 3.2, 7.2, 1.5, size=10, color=GREEN)
    
    footer(slide, "Liquidity is the fuel that moves markets")
    
    # Slide 3: BSL Visual
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Buy-Side Liquidity (BSL)")
    
    text(slide, "BSL = Stop losses ABOVE swing highs", 1.2, 0.5, 9, size=22, color=WHITE)
    
    # Visual diagram
    # Draw price action with swing high
    candle(slide, 2, 3.5, bullish=True, body=0.6, up_wick=0.1, dn_wick=0.2, scale=0.8)
    candle(slide, 2.5, 3.2, bullish=True, body=0.5, up_wick=0.1, dn_wick=0.15, scale=0.8)
    candle(slide, 3, 2.8, bullish=True, body=0.4, up_wick=0.3, dn_wick=0.1, scale=0.8)  # Swing high
    candle(slide, 3.5, 3.1, bullish=False, body=0.5, up_wick=0.15, dn_wick=0.1, scale=0.8)
    candle(slide, 4, 3.4, bullish=False, body=0.4, up_wick=0.1, dn_wick=0.2, scale=0.8)
    
    # Swing high line
    hline(slide, 2.8, 2.5, 2, color=GRAY)
    text(slide, "Swing High", 2.5, 2.5, 2, size=12, color=GRAY)
    
    # BSL zone above
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.3), Inches(2.3), Inches(2.4), Inches(0.4))
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(230, 57, 70)
    s.fill.fore_color.brightness = 0.7
    s.line.fill.background()
    text(slide, "BSL ZONE", 2.35, 2.3, 2.4, size=12, color=WHITE, bold=True)
    
    # Explanation
    text(slide, "Short sellers place stops HERE", 2.1, 5.5, 4, size=16, color=RED_ACCENT)
    text(slide, "When price sweeps this zone:", 3.0, 5.5, 4, size=16, color=WHITE)
    text(slide, "1. Buy stops trigger (become market buys)", 3.4, 5.5, 4, size=14, color=GRAY, align="left")
    text(slide, "2. Smart money sells into this buying", 3.75, 5.5, 4, size=14, color=GRAY, align="left")
    text(slide, "3. Price reverses DOWN", 4.1, 5.5, 4, size=14, color=GRAY, align="left")
    
    # Real chart placeholder
    image_placeholder(slide, 4.8, 0.5, 4, 1.8, "INSERT: BSL Sweep Chart Example")
    
    # Slide 4: SSL Visual
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Sell-Side Liquidity (SSL)")
    
    text(slide, "SSL = Stop losses BELOW swing lows", 1.2, 0.5, 9, size=22, color=WHITE)
    
    # Visual diagram - downtrend with swing low
    candle(slide, 2, 2.5, bullish=False, body=0.5, up_wick=0.2, dn_wick=0.1, scale=0.8)
    candle(slide, 2.5, 2.8, bullish=False, body=0.4, up_wick=0.15, dn_wick=0.1, scale=0.8)
    candle(slide, 3, 3.2, bullish=False, body=0.35, up_wick=0.1, dn_wick=0.4, scale=0.8)  # Swing low
    candle(slide, 3.5, 2.9, bullish=True, body=0.4, up_wick=0.1, dn_wick=0.15, scale=0.8)
    candle(slide, 4, 2.6, bullish=True, body=0.5, up_wick=0.1, dn_wick=0.2, scale=0.8)
    
    # Swing low line
    hline(slide, 3.95, 2.5, 2, color=GRAY)
    text(slide, "Swing Low", 4.05, 2.5, 2, size=12, color=GRAY)
    
    # SSL zone below
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.3), Inches(4.1), Inches(2.4), Inches(0.4))
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(34, 197, 94)
    s.fill.fore_color.brightness = 0.7
    s.line.fill.background()
    text(slide, "SSL ZONE", 4.15, 2.3, 2.4, size=12, color=WHITE, bold=True)
    
    # Explanation
    text(slide, "Long traders place stops HERE", 2.1, 5.5, 4, size=16, color=GREEN)
    text(slide, "When price sweeps this zone:", 3.0, 5.5, 4, size=16, color=WHITE)
    text(slide, "1. Sell stops trigger (become market sells)", 3.4, 5.5, 4, size=14, color=GRAY, align="left")
    text(slide, "2. Smart money buys into this selling", 3.75, 5.5, 4, size=14, color=GRAY, align="left")
    text(slide, "3. Price reverses UP", 4.1, 5.5, 4, size=14, color=GRAY, align="left")
    
    image_placeholder(slide, 4.8, 0.5, 4, 1.8, "INSERT: SSL Sweep Chart Example")
    
    # Slide 5: Equal Highs/Lows
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Equal Highs & Lows = Liquidity Magnets")
    
    # Equal highs diagram
    text(slide, "EQUAL HIGHS (EQH)", 1.2, 0.5, 4.5, size=20, color=RED_ACCENT, bold=True)
    
    # Draw 3 candles with equal highs
    candle(slide, 1.3, 2.2, bullish=True, body=0.5, up_wick=0.25, dn_wick=0.15, scale=0.7)
    candle(slide, 1.8, 2.5, bullish=False, body=0.4, up_wick=0.15, dn_wick=0.2, scale=0.7)
    candle(slide, 2.3, 2.2, bullish=True, body=0.5, up_wick=0.25, dn_wick=0.15, scale=0.7)
    
    # Equal high line
    hline(slide, 2.2, 1.0, 1.8, color=RED_ACCENT)
    text(slide, "Same Level = More Stops", 1.9, 0.8, 2.2, size=11, color=RED_ACCENT)
    
    # Equal lows diagram  
    text(slide, "EQUAL LOWS (EQL)", 1.2, 5, 4.5, size=20, color=GREEN, bold=True)
    
    candle(slide, 5.8, 2.0, bullish=False, body=0.5, up_wick=0.15, dn_wick=0.3, scale=0.7)
    candle(slide, 6.3, 1.7, bullish=True, body=0.4, up_wick=0.2, dn_wick=0.15, scale=0.7)
    candle(slide, 6.8, 2.0, bullish=False, body=0.5, up_wick=0.15, dn_wick=0.3, scale=0.7)
    
    hline(slide, 2.95, 5.5, 1.8, color=GREEN)
    text(slide, "Same Level = More Stops", 3.05, 5.3, 2.2, size=11, color=GREEN)
    
    bullets(slide, [
        "Double/Triple tops and bottoms are LIQUIDITY POOLS",
        "More touches at same level = More stops accumulated",
        "Smart money LOVES to sweep these levels",
        "Mark them on your charts!"
    ], 3.8, size=18)
    
    image_placeholder(slide, 5.5, 0.5, 9, 1.2, "INSERT: Equal Highs/Lows Sweep Example")
    
    # Slide 6: Liquidity Sweep Example
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Liquidity Sweep in Action")
    
    text(slide, "Step-by-Step Liquidity Sweep:", 1.1, 0.5, 9, size=22, color=WHITE, bold=True)
    
    # Step 1
    text(slide, "1", 1.7, 0.6, 0.4, size=20, color=RED_ACCENT, bold=True)
    text(slide, "Price creates swing high - retail goes SHORT with stops above", 1.7, 1.0, 8, size=16, color=WHITE, align="left")
    
    # Step 2
    text(slide, "2", 2.2, 0.6, 0.4, size=20, color=RED_ACCENT, bold=True)
    text(slide, "Price rallies and SWEEPS the high (takes out stops)", 2.2, 1.0, 8, size=16, color=WHITE, align="left")
    
    # Step 3
    text(slide, "3", 2.7, 0.6, 0.4, size=20, color=RED_ACCENT, bold=True)
    text(slide, "Buy stops trigger - smart money SELLS into this buying", 2.7, 1.0, 8, size=16, color=WHITE, align="left")
    
    # Step 4
    text(slide, "4", 3.2, 0.6, 0.4, size=20, color=RED_ACCENT, bold=True)
    text(slide, "Price REVERSES and drops - targeting SSL below", 3.2, 1.0, 8, size=16, color=WHITE, align="left")
    
    # Large chart placeholder
    image_placeholder(slide, 3.9, 0.5, 9, 2.7, "INSERT: Full Liquidity Sweep Trade Example")
    
    footer(slide, "The sweep IS the setup")
    
    # Slide 7: How to Trade
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "How to Trade Liquidity")
    
    bullets(slide, [
        "1. Mark key liquidity levels (highs/lows, equal H/L)",
        "2. Wait for price to SWEEP that level",
        "3. Look for reversal confirmation (rejection candle, BOS)",
        "4. Enter OPPOSITE direction of the sweep",
        "5. Target the opposite liquidity pool"
    ], 1.2, size=22)
    
    # Simple visual
    text(slide, "SWEEP HIGH", 4.5, 1, 2, size=14, color=RED_ACCENT, bold=True)
    arrow_down(slide, 2, 4.9, 0.5, RED_ACCENT)
    text(slide, "= SHORT", 5.5, 1, 2, size=14, color=RED_ACCENT)
    
    text(slide, "SWEEP LOW", 4.5, 5, 2, size=14, color=GREEN, bold=True)
    arrow_up(slide, 6, 5.4, 0.5, GREEN)
    text(slide, "= LONG", 5.0, 5, 2, size=14, color=GREEN)
    
    image_placeholder(slide, 5.8, 2.5, 5, 1, "INSERT: Entry Example After Sweep")
    
    # Slide 8: Key Takeaways
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Key Takeaways")
    
    # Two column summary with icons
    text(slide, "BSL", 1.4, 1, 1, size=28, color=RED_ACCENT, bold=True)
    text(slide, "Above highs", 1.9, 0.5, 2, size=16, color=GRAY)
    text(slide, "Target when BEARISH", 2.25, 0.5, 2, size=16, color=WHITE)
    
    text(slide, "SSL", 1.4, 4, 1, size=28, color=GREEN, bold=True)
    text(slide, "Below lows", 1.9, 3.5, 2, size=16, color=GRAY)
    text(slide, "Target when BULLISH", 2.25, 3.5, 2, size=16, color=WHITE)
    
    text(slide, "EQH/EQL", 1.4, 7, 2, size=28, color=YELLOW, bold=True)
    text(slide, "Strongest pools", 1.9, 6.5, 2.5, size=16, color=GRAY)
    text(slide, "Highest probability", 2.25, 6.5, 2.5, size=16, color=WHITE)
    
    bullets(slide, [
        "Smart money sweeps liquidity before real moves",
        "Don't place stops at obvious levels",
        "Trade WITH smart money, not against them"
    ], 3.2, size=20)
    
    text(slide, "Follow the liquidity, find the money", 5.5, 0.5, 9, size=24, color=RED_ACCENT, bold=True)
    
    prs.save("/root/clawd/dojo-gains/course-content/01_Liquidity_Visual.pptx")
    print("Saved: 01_Liquidity_Visual.pptx")


# ============================================
# FVG/IMBALANCES - WITH VISUALS
# ============================================
def create_fvg_visual():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    text(slide, "IMBALANCES", 2.2, 0.5, 9, size=52, color=WHITE, bold=True)
    text(slide, "Fair Value Gaps (FVG)", 3.1, 0.5, 9, size=26, color=GRAY)
    hline(slide, 3.9, 3.5, 3)
    text(slide, "DOJO GAINS - Core Pillar 3", 4.4, 0.5, 9, size=20, color=RED_ACCENT, bold=True)
    
    # Slide 2: What is FVG - Visual
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "What is a Fair Value Gap?")
    
    text(slide, "An FVG is a GAP created when price moves too fast", 1.1, 0.5, 9, size=20, color=WHITE)
    
    # Visual: 3 candles with gap
    text(slide, "3-CANDLE FORMATION:", 1.7, 0.5, 4, size=16, color=GRAY)
    
    # Candle 1
    candle(slide, 1.5, 2.8, bullish=True, body=0.5, up_wick=0.15, dn_wick=0.15, scale=1.0)
    text(slide, "Candle 1", 3.8, 1.1, 1, size=12, color=GRAY)
    
    # Candle 2 (displacement)
    candle(slide, 2.3, 2.0, bullish=True, body=1.0, up_wick=0.1, dn_wick=0.1, scale=1.0)
    text(slide, "Candle 2", 3.4, 1.9, 1.2, size=12, color=GREEN)
    text(slide, "(Displacement)", 3.65, 1.8, 1.4, size=10, color=GREEN)
    
    # Candle 3
    candle(slide, 3.1, 1.6, bullish=True, body=0.5, up_wick=0.15, dn_wick=0.15, scale=1.0)
    text(slide, "Candle 3", 2.6, 2.7, 1, size=12, color=GRAY)
    
    # FVG zone highlight
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(2.55), Inches(2.5), Inches(0.35))
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(59, 130, 246)
    s.fill.fore_color.brightness = 0.6
    s.line.fill.background()
    text(slide, "FVG", 2.55, 1.0, 2.5, size=14, color=WHITE, bold=True)
    
    # Labels
    text(slide, "C1 High", 2.7, 0.3, 0.8, size=10, color=GRAY)
    text(slide, "C3 Low", 2.35, 3.4, 0.8, size=10, color=GRAY)
    
    # Explanation
    text(slide, "The GAP between:", 1.7, 5, 4.5, size=18, color=WHITE, bold=True)
    text(slide, "Candle 1's HIGH and Candle 3's LOW", 2.1, 5, 4.5, size=16, color=BLUE)
    text(slide, "This gap = IMBALANCE = Entry zone", 2.5, 5, 4.5, size=16, color=GRAY)
    text(slide, "Price often RETURNS to fill this gap", 2.9, 5, 4.5, size=16, color=GREEN)
    
    image_placeholder(slide, 4.0, 5, 4.5, 2.5, "INSERT: FVG Example on Chart")
    
    # Slide 3: Bullish FVG
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Bullish Fair Value Gap")
    
    # Visual
    candle(slide, 1.8, 3.2, bullish=True, body=0.4, up_wick=0.1, dn_wick=0.2, scale=1.0)
    candle(slide, 2.5, 2.4, bullish=True, body=0.9, up_wick=0.08, dn_wick=0.08, scale=1.0)
    candle(slide, 3.2, 1.8, bullish=True, body=0.5, up_wick=0.15, dn_wick=0.12, scale=1.0)
    
    # FVG zone
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.3), Inches(2.8), Inches(2.4), Inches(0.45))
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(34, 197, 94)
    s.fill.fore_color.brightness = 0.6
    s.line.fill.background()
    text(slide, "BULLISH FVG", 2.85, 1.3, 2.4, size=12, color=WHITE, bold=True)
    
    # Arrow showing price return
    arrow_down(slide, 4.2, 1.5, 1.2, GRAY)
    text(slide, "Price returns", 1.3, 4.5, 1.5, size=11, color=GRAY)
    text(slide, "to FVG", 1.55, 4.5, 1.5, size=11, color=GRAY)
    
    bullets(slide, [
        "Forms during UPWARD moves",
        "Candle 3 LOW > Candle 1 HIGH = Gap",
        "Shows aggressive BUYING",
        "FVG acts as SUPPORT",
        "Enter LONG when price retraces to FVG"
    ], 1.2, size=20)
    
    image_placeholder(slide, 4.5, 0.5, 4.5, 2, "INSERT: Bullish FVG Trade Example")
    
    text(slide, "Bullish FVG = BUY ZONE", 4.7, 5.2, 4.3, size=22, color=GREEN, bold=True)
    
    # Slide 4: Bearish FVG
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Bearish Fair Value Gap")
    
    # Visual
    candle(slide, 1.8, 1.8, bullish=False, body=0.4, up_wick=0.2, dn_wick=0.1, scale=1.0)
    candle(slide, 2.5, 2.4, bullish=False, body=0.9, up_wick=0.08, dn_wick=0.08, scale=1.0)
    candle(slide, 3.2, 3.0, bullish=False, body=0.5, up_wick=0.12, dn_wick=0.15, scale=1.0)
    
    # FVG zone
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.3), Inches(2.25), Inches(2.4), Inches(0.45))
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(230, 57, 70)
    s.fill.fore_color.brightness = 0.6
    s.line.fill.background()
    text(slide, "BEARISH FVG", 2.3, 1.3, 2.4, size=12, color=WHITE, bold=True)
    
    # Arrow showing price return
    arrow_up(slide, 4.2, 2.4, 1.2, GRAY)
    text(slide, "Price returns", 3.6, 4.5, 1.5, size=11, color=GRAY)
    text(slide, "to FVG", 3.85, 4.5, 1.5, size=11, color=GRAY)
    
    bullets(slide, [
        "Forms during DOWNWARD moves",
        "Candle 3 HIGH < Candle 1 LOW = Gap",
        "Shows aggressive SELLING",
        "FVG acts as RESISTANCE",
        "Enter SHORT when price retraces to FVG"
    ], 1.2, size=20)
    
    image_placeholder(slide, 4.5, 0.5, 4.5, 2, "INSERT: Bearish FVG Trade Example")
    
    text(slide, "Bearish FVG = SELL ZONE", 4.7, 5.2, 4.3, size=22, color=RED_ACCENT, bold=True)
    
    # Slide 5: 50% CE Level
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "The 50% Level (Consequent Encroachment)")
    
    # Visual of FVG with 50% marked
    candle(slide, 2, 2.8, bullish=True, body=0.4, up_wick=0.1, dn_wick=0.15, scale=1.0)
    candle(slide, 2.8, 2.0, bullish=True, body=0.95, up_wick=0.08, dn_wick=0.08, scale=1.0)
    candle(slide, 3.6, 1.5, bullish=True, body=0.5, up_wick=0.15, dn_wick=0.1, scale=1.0)
    
    # FVG zone with 50% line
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(2.55), Inches(2.6), Inches(0.55))
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(59, 130, 246)
    s.fill.fore_color.brightness = 0.65
    s.line.fill.background()
    
    # 50% line
    hline(slide, 2.82, 1.5, 2.6, color=YELLOW)
    text(slide, "50% (CE)", 2.65, 4.2, 1.5, size=14, color=YELLOW, bold=True)
    
    text(slide, "FVG Top", 2.45, 4.2, 1.2, size=11, color=GRAY)
    text(slide, "FVG Bottom", 3.0, 4.2, 1.2, size=11, color=GRAY)
    
    bullets(slide, [
        "CE = 50% midpoint of the FVG",
        "This is the EQUILIBRIUM of the imbalance",
        "BEST entry level within the FVG",
        "Price often reacts precisely at CE",
        "Use CE as your primary entry point"
    ], 1.2, size=20)
    
    text(slide, "50% CE = Sweet Spot Entry", 5.5, 0.5, 9, size=24, color=YELLOW, bold=True)
    
    # Slide 6: Trading FVGs
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "How to Trade FVGs")
    
    text(slide, "Step-by-Step FVG Trade:", 1.1, 0.5, 9, size=22, color=WHITE, bold=True)
    
    text(slide, "1", 1.7, 0.6, 0.4, size=20, color=BLUE, bold=True)
    text(slide, "Identify the TREND direction on higher timeframe", 1.7, 1.0, 8, size=16, color=WHITE, align="left")
    
    text(slide, "2", 2.2, 0.6, 0.4, size=20, color=BLUE, bold=True)
    text(slide, "Look for DISPLACEMENT candle creating FVG", 2.2, 1.0, 8, size=16, color=WHITE, align="left")
    
    text(slide, "3", 2.7, 0.6, 0.4, size=20, color=BLUE, bold=True)
    text(slide, "Mark the FVG zone and 50% CE level", 2.7, 1.0, 8, size=16, color=WHITE, align="left")
    
    text(slide, "4", 3.2, 0.6, 0.4, size=20, color=BLUE, bold=True)
    text(slide, "WAIT for price to retrace back to FVG", 3.2, 1.0, 8, size=16, color=WHITE, align="left")
    
    text(slide, "5", 3.7, 0.6, 0.4, size=20, color=BLUE, bold=True)
    text(slide, "Enter at 50% CE with confirmation candle", 3.7, 1.0, 8, size=16, color=WHITE, align="left")
    
    text(slide, "6", 4.2, 0.6, 0.4, size=20, color=BLUE, bold=True)
    text(slide, "Stop beyond FVG, target next liquidity", 4.2, 1.0, 8, size=16, color=WHITE, align="left")
    
    image_placeholder(slide, 4.9, 0.5, 9, 1.7, "INSERT: Complete FVG Trade Example (Entry to Exit)")
    
    # Slide 7: Key Takeaways
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Key Takeaways")
    
    bullets(slide, [
        "FVG = Gap where price moved too fast",
        "3-candle formation: Check C1 & C3 wicks",
        "Bullish FVG = Support (buy zone)",
        "Bearish FVG = Resistance (sell zone)",
        "50% CE is the optimal entry level",
        "Wait for price to RETURN to FVG before entering"
    ], 1.2, size=22)
    
    text(slide, "Imbalances are magnets - price comes back", 5.5, 0.5, 9, size=24, color=RED_ACCENT, bold=True)
    
    prs.save("/root/clawd/dojo-gains/course-content/03_Imbalances_Visual.pptx")
    print("Saved: 03_Imbalances_Visual.pptx")


# ============================================
# FIBONACCI - WITH VISUALS
# ============================================
def create_fib_visual():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    text(slide, "FIBONACCI", 2.2, 0.5, 9, size=52, color=WHITE, bold=True)
    text(slide, "The Golden Ratio in Trading", 3.1, 0.5, 9, size=26, color=GRAY)
    hline(slide, 3.9, 3.5, 3)
    text(slide, "DOJO GAINS - Mentorship", 4.4, 0.5, 9, size=20, color=RED_ACCENT, bold=True)
    
    # Slide 2: Key Levels Visual
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Key Fibonacci Levels")
    
    # Visual representation of fib levels
    # Draw vertical bar showing levels
    base_x = 1.5
    base_top = 1.4
    base_height = 4.5
    bar_width = 0.8
    
    # Full bar
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
        Inches(base_x), Inches(base_top), Inches(bar_width), Inches(base_height))
    s.fill.solid()
    s.fill.fore_color.rgb = DARK_GRAY
    s.line.fill.background()
    
    # Level lines and labels
    levels = [
        (0, "0% (Swing Low)", GRAY),
        (0.236, "23.6%", GRAY),
        (0.382, "38.2%", BLUE),
        (0.5, "50%", WHITE),
        (0.618, "61.8% (Golden)", YELLOW),
        (0.786, "78.6%", RED_ACCENT),
        (1.0, "100% (Swing High)", GRAY),
    ]
    
    for level, label, color in levels:
        y = base_top + base_height * (1 - level)
        hline(slide, y, base_x, bar_width + 0.5, color)
        text(slide, label, y - 0.15, base_x + bar_width + 0.6, 2.5, size=14, color=color, align="left")
    
    # OTE Zone highlight
    ote_top = base_top + base_height * (1 - 0.786)
    ote_height = base_height * (0.786 - 0.618)
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(base_x), Inches(ote_top), Inches(bar_width), Inches(ote_height))
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(34, 197, 94)
    s.fill.fore_color.brightness = 0.5
    s.line.fill.background()
    text(slide, "OTE ZONE", ote_top + 0.1, base_x + 0.1, 0.7, size=10, color=WHITE, bold=True)
    
    # Explanation on right
    text(slide, "KEY LEVELS:", 1.3, 5, 4.5, size=20, color=WHITE, bold=True)
    bullets(slide, [
        "61.8% = Golden Ratio (most important)",
        "50% = Equilibrium",
        "OTE Zone (62-79%) = Best entries",
        "38.2% = Shallow retracement"
    ], 1.8, size=18)
    
    image_placeholder(slide, 4.5, 5, 4.5, 2, "INSERT: Fib Retracement on Chart")
    
    # Slide 3: Premium vs Discount
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Premium vs Discount Zones")
    
    # Visual
    base_x = 2.5
    
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(base_x), Inches(1.5), Inches(1), Inches(2.2))
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(230, 57, 70)
    s.fill.fore_color.brightness = 0.5
    s.line.fill.background()
    text(slide, "PREMIUM", 2.3, base_x, 1, size=14, color=WHITE, bold=True)
    text(slide, "Above 50%", 2.6, base_x, 1, size=11, color=WHITE)
    text(slide, "SELL ZONE", 2.9, base_x, 1, size=12, color=RED_ACCENT, bold=True)
    
    hline(slide, 3.7, base_x - 0.3, 1.6, WHITE)
    text(slide, "50%", 3.55, base_x + 1.1, 0.6, size=12, color=WHITE, bold=True)
    
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(base_x), Inches(3.75), Inches(1), Inches(2.2))
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(34, 197, 94)
    s.fill.fore_color.brightness = 0.5
    s.line.fill.background()
    text(slide, "DISCOUNT", 4.5, base_x, 1, size=14, color=WHITE, bold=True)
    text(slide, "Below 50%", 4.8, base_x, 1, size=11, color=WHITE)
    text(slide, "BUY ZONE", 5.1, base_x, 1, size=12, color=GREEN, bold=True)
    
    # Rules on right
    text(slide, "THE RULE:", 1.5, 5, 4.5, size=22, color=WHITE, bold=True)
    
    text(slide, "In UPTREND:", 2.2, 5, 4.5, size=18, color=GREEN, bold=True)
    text(slide, "Buy at DISCOUNT (below 50%)", 2.6, 5, 4.5, size=16, color=WHITE)
    
    text(slide, "In DOWNTREND:", 3.3, 5, 4.5, size=18, color=RED_ACCENT, bold=True)
    text(slide, "Sell at PREMIUM (above 50%)", 3.7, 5, 4.5, size=16, color=WHITE)
    
    text(slide, "Buy cheap, sell expensive", 5.5, 0.5, 9, size=24, color=YELLOW, bold=True)
    
    # Slide 4: OTE Zone
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "OTE - Optimal Trade Entry")
    
    bullets(slide, [
        "OTE = 62% to 79% Fibonacci zone",
        "This is where smart money enters",
        "Deepest pullback before continuation",
        "Combine with FVG/OB for A+ setup"
    ], 1.2, size=22)
    
    # Visual
    text(slide, "OTE ZONE VISUAL:", 3.5, 0.5, 4, size=16, color=GRAY)
    
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(1), Inches(4), Inches(3.5), Inches(1.2))
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(34, 197, 94)
    s.fill.fore_color.brightness = 0.5
    s.line.fill.background()
    
    text(slide, "62%", 3.85, 1, 0.8, size=12, color=WHITE)
    text(slide, "OTE ZONE - ENTER HERE", 4.35, 1, 3.5, size=16, color=WHITE, bold=True)
    text(slide, "79%", 5.0, 1, 0.8, size=12, color=WHITE)
    
    image_placeholder(slide, 3.8, 5, 4.5, 2.8, "INSERT: OTE Entry Example")
    
    text(slide, "OTE + FVG + Killzone = Sniper Entry", 6.5, 0.5, 9, size=20, color=GREEN, bold=True)
    
    # Slide 5: Key Takeaways
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Key Takeaways")
    
    bullets(slide, [
        "61.8% = Golden Ratio (most important)",
        "OTE Zone (62-79%) = Optimal entry area",
        "Premium = Sell zone | Discount = Buy zone",
        "Combine Fib with FVG/OB for best entries",
        "Draw Fib from swing low to high (or vice versa)",
        "Wait for price to enter OTE before looking for entry"
    ], 1.2, size=22)
    
    image_placeholder(slide, 4.5, 0.5, 9, 2, "INSERT: Complete Fib Trade Example")
    
    prs.save("/root/clawd/dojo-gains/course-content/M06_Fibonacci_Visual.pptx")
    print("Saved: M06_Fibonacci_Visual.pptx")


if __name__ == "__main__":
    create_liquidity_visual()
    create_fvg_visual()
    create_fib_visual()
    print("\nVisual presentations created!")
    print("NOTE: Add your own chart screenshots where indicated by placeholders")
