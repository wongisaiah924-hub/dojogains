#!/usr/bin/env python3
"""
Dojo Gains - Candlestick Anatomy & Market Structure Presentation
v8 - Fixed all spacing issues
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
DARK_BG = RGBColor(15, 15, 26)
RED_ACCENT = RGBColor(230, 57, 70)
GREEN = RGBColor(34, 197, 94)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(156, 163, 175)

SLIDE_W = 10
SLIDE_H = 7.5

def set_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK_BG

def text(slide, txt, top, left, width, size=18, color=WHITE, bold=False, align="center"):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = PP_ALIGN.CENTER if align == "center" else PP_ALIGN.LEFT
    return box

def title(slide, txt, top=0.4):
    text(slide, txt, top, 0.5, 9, size=40, color=WHITE, bold=True)

def bullets(slide, points, top, size=20):
    box = slide.shapes.add_textbox(Inches(0.75), Inches(top), Inches(8.5), Inches(5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, pt in enumerate(points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•   {pt}"
        p.font.size = Pt(size)
        p.font.color.rgb = WHITE
        p.space_before = Pt(6)
        p.space_after = Pt(10)

def footer(slide, txt):
    text(slide, txt, 6.9, 0.5, 9, size=14, color=GRAY)

def candle_bottom_aligned(slide, cx, baseline, bullish=True, body=1.0, up_wick=0.3, dn_wick=0.3, scale=1.0):
    """Draw candlestick with BOTTOM at baseline"""
    color = GREEN if bullish else RED_ACCENT
    bw = 0.45 * scale
    ww = 0.06 * scale
    body *= scale
    up_wick *= scale
    dn_wick *= scale
    total_height = up_wick + body + dn_wick
    top = baseline - total_height
    
    if up_wick > 0.01:
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
            Inches(cx - ww/2), Inches(top), Inches(ww), Inches(up_wick))
        s.fill.solid()
        s.fill.fore_color.rgb = color
        s.line.fill.background()
    
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(cx - bw/2), Inches(top + up_wick), Inches(bw), Inches(body))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    
    if dn_wick > 0.01:
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(cx - ww/2), Inches(top + up_wick + body), Inches(ww), Inches(dn_wick))
        s.fill.solid()
        s.fill.fore_color.rgb = color
        s.line.fill.background()

def candle(slide, cx, top, bullish=True, body=1.0, up_wick=0.3, dn_wick=0.3, scale=1.0):
    """Draw candlestick from TOP position"""
    color = GREEN if bullish else RED_ACCENT
    bw = 0.45 * scale
    ww = 0.06 * scale
    body *= scale
    up_wick *= scale
    dn_wick *= scale
    
    if up_wick > 0.01:
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
            Inches(cx - ww/2), Inches(top), Inches(ww), Inches(up_wick))
        s.fill.solid()
        s.fill.fore_color.rgb = color
        s.line.fill.background()
    
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(cx - bw/2), Inches(top + up_wick), Inches(bw), Inches(body))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    
    if dn_wick > 0.01:
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(cx - ww/2), Inches(top + up_wick + body), Inches(ww), Inches(dn_wick))
        s.fill.solid()
        s.fill.fore_color.rgb = color
        s.line.fill.background()

def hline(slide, top, left, width, color=RED_ACCENT):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.02))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    
    # ===== SLIDE 1: Title =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    text(slide, "JAPANESE CANDLESTICKS", 2.2, 0.5, 9, size=48, color=WHITE, bold=True)
    text(slide, "Anatomy & Market Structure", 3.0, 0.5, 9, size=28, color=GRAY)
    hline(slide, 3.8, 3.5, 3)
    text(slide, "DOJO GAINS", 4.3, 0.5, 9, size=24, color=RED_ACCENT, bold=True)
    text(slide, "Master the Art of Trading", 4.9, 0.5, 9, size=16, color=GRAY)
    
    # ===== SLIDE 2: What Are Candlesticks =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "What Are Candlesticks?")
    bullets(slide, [
        "Originated in 18th century Japan by rice trader Munehisa Homma",
        "Each candle represents price action over a specific time period",
        "Shows FOUR key prices: Open, High, Low, Close (OHLC)",
        "More visual than line charts - tells the story of price",
        "The foundation of technical analysis and price action trading",
        "Used across all markets: stocks, options, crypto, futures"
    ], 1.2, size=22)
    footer(slide, "Candlesticks reveal market psychology in a single glance")
    
    # ===== SLIDE 3: Anatomy =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Candlestick Anatomy")
    
    # Bullish candle at x=3, starts at y=1.5
    # up_wick=0.5 (1.5-2.0), body=1.5 (2.0-3.5), dn_wick=0.4 (3.5-3.9)
    candle(slide, 3, 1.5, bullish=True, body=1.5, up_wick=0.5, dn_wick=0.4, scale=1.0)
    text(slide, "HIGH", 1.3, 0.8, 1.5, size=14, color=GRAY)
    text(slide, "CLOSE", 2.0, 0.8, 1.5, size=14, color=GREEN, bold=True)
    text(slide, "OPEN", 3.2, 0.8, 1.5, size=14, color=GREEN, bold=True)
    text(slide, "LOW", 3.9, 0.8, 1.5, size=14, color=GRAY)
    text(slide, "BULLISH", 4.6, 2.3, 1.5, size=20, color=GREEN, bold=True)
    
    # Bearish candle at x=7
    candle(slide, 7, 1.5, bullish=False, body=1.5, up_wick=0.4, dn_wick=0.5, scale=1.0)
    text(slide, "HIGH", 1.3, 7.7, 1.5, size=14, color=GRAY)
    text(slide, "OPEN", 2.0, 7.7, 1.5, size=14, color=RED_ACCENT, bold=True)
    text(slide, "CLOSE", 3.2, 7.7, 1.5, size=14, color=RED_ACCENT, bold=True)
    text(slide, "LOW", 4.0, 7.7, 1.5, size=14, color=GRAY)
    text(slide, "BEARISH", 4.6, 6.3, 1.5, size=20, color=RED_ACCENT, bold=True)
    
    text(slide, "Upper Wick = Price reached higher but pulled back", 5.5, 0.5, 9, size=16, color=GRAY)
    text(slide, "Body = Open to Close range", 5.95, 0.5, 9, size=16, color=GRAY)
    text(slide, "Lower Wick = Price went lower but recovered", 6.4, 0.5, 9, size=16, color=GRAY)
    
    # ===== SLIDE 4: Bullish vs Bearish =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Bullish vs Bearish")
    
    # Left: Bullish
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), Inches(1.35), Inches(0.35), Inches(0.35))
    circle.fill.solid()
    circle.fill.fore_color.rgb = GREEN
    circle.line.fill.background()
    text(slide, "BULLISH", 1.25, 1.3, 3.5, size=28, color=GREEN, bold=True, align="left")
    candle(slide, 2.75, 2.0, bullish=True, body=1.3, up_wick=0.25, dn_wick=0.2, scale=1.0)
    text(slide, "Close > Open", 4.0, 0.5, 4.5, size=20, color=WHITE)
    text(slide, "Buyers won", 4.5, 0.5, 4.5, size=20, color=WHITE)
    text(slide, "Price moved UP", 5.0, 0.5, 4.5, size=20, color=WHITE)
    
    # Right: Bearish
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.3), Inches(1.35), Inches(0.35), Inches(0.35))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RED_ACCENT
    circle.line.fill.background()
    text(slide, "BEARISH", 1.25, 5.8, 3.5, size=28, color=RED_ACCENT, bold=True, align="left")
    candle(slide, 7.25, 2.0, bullish=False, body=1.3, up_wick=0.2, dn_wick=0.25, scale=1.0)
    text(slide, "Close < Open", 4.0, 5, 4.5, size=20, color=WHITE)
    text(slide, "Sellers won", 4.5, 5, 4.5, size=20, color=WHITE)
    text(slide, "Price moved DOWN", 5.0, 5, 4.5, size=20, color=WHITE)
    
    footer(slide, "The color tells you WHO WON - buyers or sellers")
    
    # ===== SLIDE 5: Candle Size =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Candle Size = Momentum")
    
    baseline = 3.3
    
    # Small (center at 2)
    text(slide, "SMALL", 1.1, 0.5, 3, size=20, color=GRAY, bold=True)
    candle_bottom_aligned(slide, 2, baseline, bullish=True, body=0.35, up_wick=0.12, dn_wick=0.12, scale=1.0)
    text(slide, "Indecision", 3.5, 0.5, 3, size=16, color=GRAY)
    text(slide, "Low momentum", 3.9, 0.5, 3, size=16, color=GRAY)
    
    # Medium (center at 5)
    text(slide, "MEDIUM", 1.1, 3.5, 3, size=20, color=WHITE, bold=True)
    candle_bottom_aligned(slide, 5, baseline, bullish=True, body=0.65, up_wick=0.15, dn_wick=0.12, scale=1.0)
    text(slide, "Normal", 3.5, 3.5, 3, size=16, color=WHITE)
    text(slide, "Steady movement", 3.9, 3.5, 3, size=16, color=WHITE)
    
    # Large (center at 8)
    text(slide, "LARGE", 1.1, 6.5, 3, size=20, color=GREEN, bold=True)
    candle_bottom_aligned(slide, 8, baseline, bullish=True, body=1.1, up_wick=0.08, dn_wick=0.06, scale=1.0)
    text(slide, "Strong conviction", 3.5, 6.5, 3, size=16, color=GREEN)
    text(slide, "High momentum!", 3.9, 6.5, 3, size=16, color=GREEN)
    
    bullets(slide, [
        "Big green candle = Strong buying pressure",
        "Big red candle = Strong selling pressure",
        "Small candles = Market undecided"
    ], 4.6, size=20)
    
    # ===== SLIDE 6: Wicks =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Reading the Wicks")
    
    bullets(slide, [
        "Wicks show REJECTION - price tried but got pushed back",
        "Long upper wick = Sellers rejected higher prices",
        "Long lower wick = Buyers rejected lower prices",
        "No wick = Strong conviction in that direction"
    ], 1.2, size=22)
    
    baseline = 5.8
    
    # Bearish rejection
    candle_bottom_aligned(slide, 2, baseline, bullish=False, body=0.35, up_wick=0.9, dn_wick=0.08, scale=1.0)
    text(slide, "Bearish Rejection", 6.0, 0.5, 3, size=16, color=RED_ACCENT, bold=True)
    
    # Bullish rejection
    candle_bottom_aligned(slide, 5, baseline, bullish=True, body=0.35, up_wick=0.08, dn_wick=0.9, scale=1.0)
    text(slide, "Bullish Rejection", 6.0, 3.5, 3, size=16, color=GREEN, bold=True)
    
    # Full conviction
    candle_bottom_aligned(slide, 8, baseline, bullish=True, body=1.0, up_wick=0.0, dn_wick=0.0, scale=1.0)
    text(slide, "Full Conviction", 6.0, 6.5, 3, size=16, color=WHITE, bold=True)
    
    # ===== SLIDE 7: Single Patterns =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Single Candle Patterns")
    
    baseline = 3.8
    cols = [1.5, 3.8, 6.1, 8.4]
    col_w = 2.0
    
    patterns = [
        ("DOJI", "Indecision", "Reversal signal", True, 0.06, 0.5, 0.5),
        ("HAMMER", "Bullish reversal", "At support", True, 0.35, 0.06, 0.75),
        ("SHOOTING STAR", "Bearish reversal", "At resistance", False, 0.35, 0.75, 0.06),
        ("MARUBOZU", "Full conviction", "Strong trend", True, 1.0, 0.0, 0.0),
    ]
    
    for i, (name, line1, line2, bull, body, up, dn) in enumerate(patterns):
        cx = cols[i]
        text(slide, name, 1.1, cx - 1.0, col_w, size=15, color=WHITE, bold=True)
        candle_bottom_aligned(slide, cx, baseline, bullish=bull, body=body, up_wick=up, dn_wick=dn, scale=1.1)
        text(slide, line1, 4.0, cx - 1.0, col_w, size=13, color=GRAY)
        text(slide, line2, 4.35, cx - 1.0, col_w, size=13, color=GRAY)
    
    footer(slide, "These patterns work best at key support/resistance levels")
    
    # ===== SLIDE 8: Multi-Candle Patterns =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Multi-Candle Patterns")
    
    baseline = 3.4
    
    # Bullish Engulfing
    text(slide, "BULLISH ENGULFING", 1.1, 0.3, 3.0, size=17, color=GREEN, bold=True)
    candle_bottom_aligned(slide, 1.3, baseline, bullish=False, body=0.45, up_wick=0.08, dn_wick=0.08, scale=1.0)
    candle_bottom_aligned(slide, 2.0, baseline, bullish=True, body=0.9, up_wick=0.1, dn_wick=0.12, scale=1.0)
    text(slide, "Green engulfs red", 3.6, 0.3, 3.0, size=13, color=GRAY)
    text(slide, "Bullish Reversal", 3.95, 0.3, 3.0, size=15, color=GREEN, bold=True)
    
    # Bearish Engulfing
    text(slide, "BEARISH ENGULFING", 1.1, 3.5, 3.0, size=17, color=RED_ACCENT, bold=True)
    candle_bottom_aligned(slide, 4.5, baseline, bullish=True, body=0.45, up_wick=0.08, dn_wick=0.08, scale=1.0)
    candle_bottom_aligned(slide, 5.2, baseline, bullish=False, body=0.9, up_wick=0.12, dn_wick=0.1, scale=1.0)
    text(slide, "Red engulfs green", 3.6, 3.5, 3.0, size=13, color=GRAY)
    text(slide, "Bearish Reversal", 3.95, 3.5, 3.0, size=15, color=RED_ACCENT, bold=True)
    
    # Inside Bar
    text(slide, "INSIDE BAR", 1.1, 6.7, 3.0, size=17, color=WHITE, bold=True)
    candle_bottom_aligned(slide, 7.7, baseline, bullish=True, body=0.9, up_wick=0.2, dn_wick=0.2, scale=1.0)
    candle_bottom_aligned(slide, 8.4, baseline, bullish=False, body=0.35, up_wick=0.08, dn_wick=0.08, scale=1.0)
    text(slide, "2nd inside 1st range", 3.6, 6.7, 3.0, size=13, color=GRAY)
    text(slide, "Breakout Setup", 3.95, 6.7, 3.0, size=15, color=WHITE, bold=True)
    
    footer(slide, "These patterns show a SHIFT in control between buyers and sellers")
    
    # ===== SLIDE 9: Trends =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Market Structure: Trends")
    
    bullets(slide, [
        "UPTREND: Higher Highs (HH) and Higher Lows (HL)",
        "DOWNTREND: Lower Highs (LH) and Lower Lows (LL)",
        "SIDEWAYS: Price bouncing between support and resistance",
        "Trade WITH the trend - it's your friend"
    ], 1.2, size=22)
    
    text(slide, "UPTREND", 4.6, 0.5, 2.8, size=22, color=GREEN, bold=True)
    text(slide, "HH > HL > HH > HL", 5.1, 0.5, 2.8, size=14, color=GRAY)
    
    text(slide, "DOWNTREND", 4.6, 3.6, 2.8, size=22, color=RED_ACCENT, bold=True)
    text(slide, "LH > LL > LH > LL", 5.1, 3.6, 2.8, size=14, color=GRAY)
    
    text(slide, "SIDEWAYS", 4.6, 6.7, 2.8, size=22, color=WHITE, bold=True)
    text(slide, "Range-bound", 5.1, 6.7, 2.8, size=14, color=GRAY)
    
    # ===== SLIDE 10: Support & Resistance =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Support & Resistance")
    
    bullets(slide, [
        "SUPPORT: Level where buying pressure stops decline",
        "RESISTANCE: Level where selling pressure stops advance",
        "Think of them as ZONES, not exact lines",
        "The more times tested, the stronger the level",
        "When broken, support becomes resistance (flip)",
        "Candle patterns AT these levels = best setups"
    ], 1.2, size=22)
    
    footer(slide, "Candlestick patterns are most reliable at major S/R levels!")
    
    # ===== SLIDE 11: Timeframes =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Timeframes Matter")
    
    bullets(slide, [
        "Each candle = one period of that timeframe",
        "1-min chart: Each candle = 1 minute of price action",
        "Daily chart: Each candle = entire trading day",
        "Higher timeframes = more reliable (less noise)",
        "Always check higher TF for context"
    ], 1.2, size=22)
    
    text(slide, "DAY TRADERS", 4.8, 1.0, 3.5, size=20, color=WHITE, bold=True)
    text(slide, "1m - 5m - 15m", 5.3, 1.0, 3.5, size=18, color=GRAY)
    
    text(slide, "SWING TRADERS", 4.8, 5.5, 3.5, size=20, color=WHITE, bold=True)
    text(slide, "1H - 4H - Daily", 5.3, 5.5, 3.5, size=18, color=GRAY)
    
    # ===== SLIDE 12: Putting It Together =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Putting It All Together")
    
    bullets(slide, [
        "1.  Identify the TREND",
        "2.  Mark key SUPPORT and RESISTANCE",
        "3.  Wait for price to reach a key level",
        "4.  Look for CANDLESTICK PATTERNS there",
        "5.  Confirm with candle SIZE and WICKS",
        "6.  Enter with the trend, manage risk"
    ], 1.2, size=24)
    
    text(slide, "Patience + Pattern + Level = High Probability", 5.3, 0.5, 9, size=26, color=RED_ACCENT, bold=True)
    
    # ===== SLIDE 13: Key Takeaways =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Key Takeaways")
    
    bullets(slide, [
        "Candlesticks = battle between buyers and sellers",
        "Body size = momentum  |  Wicks = rejection",
        "Green = bulls won  |  Red = bears won",
        "Key patterns: Doji, Hammer, Shooting Star, Engulfing",
        "Trade WITH the trend",
        "Patterns work best at S/R levels",
        "Higher timeframes = more reliable"
    ], 1.2, size=22)
    
    footer(slide, "Practice daily. Screen time builds intuition.")
    
    # ===== SLIDE 14: Next Steps =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title(slide, "Your Next Steps")
    
    bullets(slide, [
        "Open TradingView and practice identifying patterns",
        "Mark support & resistance on your favorite tickers",
        "Start a trading journal - note what you see",
        "Paper trade before risking real money",
        "Join the Dojo Discord to discuss setups"
    ], 1.2, size=24)
    
    text(slide, "dojogains.com", 5.0, 0.5, 9, size=36, color=RED_ACCENT, bold=True)
    text(slide, "No hype. No gurus. Just growth.", 5.7, 0.5, 9, size=20, color=GRAY)
    
    # Save
    path = "/root/clawd/dojo-gains/course-content/Candlestick_Anatomy_DojoGains.pptx"
    prs.save(path)
    print(f"Saved: {path}")

if __name__ == "__main__":
    create_presentation()
